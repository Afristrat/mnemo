# -*- coding: utf-8 -*-
"""
Génère la présentation « Mnémo — Modèle économique & monétisation ».
Cas d'usage pour la cohorte d'entrepreneurs élite IA (Le Labo IA).
Sortie : Mnemo-modele-economique.pptx (16:9, éditable).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Palette — sobre, un accent teal (statut « souverain », cf. système de design)
# ---------------------------------------------------------------------------
TEAL       = RGBColor(0x0E, 0x7C, 0x7B)   # accent primaire
TEAL_DARK  = RGBColor(0x0A, 0x5C, 0x5B)
TEAL_SOFT  = RGBColor(0xE3, 0xF0, 0xEF)   # fond doux
INK        = RGBColor(0x14, 0x1B, 0x1F)   # texte titre
SLATE      = RGBColor(0x3C, 0x4A, 0x50)   # texte corps
MUTE       = RGBColor(0x6B, 0x7A, 0x80)   # texte secondaire
LINE       = RGBColor(0xD7, 0xDE, 0xE0)   # filets
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PAPER      = RGBColor(0xF7, 0xF9, 0xF9)   # fond de slide
GOLD       = RGBColor(0xB8, 0x8A, 0x2E)   # intelligence avancée (parcimonie)
BLUE       = RGBColor(0x2A, 0x5C, 0x9E)   # réseau / cloud
RED_SOFT   = RGBColor(0xB0, 0x47, 0x3A)   # faiblesse / avant

EMU = 914400
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

FONT_TITLE = "Space Grotesk"   # titres (avec repli système)
FONT_BODY  = "Inter"           # corps
FONT_MONO  = "Consolas"        # données / chiffres


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    _send_back(s, bg)
    return s


def _send_back(s, shp):
    spTree = s.shapes._spTree
    spTree.remove(shp._element)
    spTree.insert(2, shp._element)


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def _set_run(r, text, size, color, bold=False, font=FONT_BODY, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def txt(s, x, y, w, h, text, size=14, color=SLATE, bold=False, font=FONT_BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.08,
        space_after=2):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    _set_run(p.add_run(), text, size, color, bold, font, italic)
    return tb


def multi(s, x, y, w, h, lines, size=14, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
          line_spacing=1.12):
    """lines = list of (text, color, bold, font, size_override, space_after)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        text, color, bold, font, so, sa = (ln + (None,) * 6)[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(sa if sa is not None else 4)
        _set_run(p.add_run(), text, so or size, color or SLATE, bool(bold),
                 font or FONT_BODY)
    return tb


def bullets(s, x, y, w, h, items, size=14, color=SLATE, gap=7, dot_color=TEAL,
            lead_bold_split=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        p.space_after = Pt(gap)
        rd = p.add_run(); _set_run(rd, "—  ", size, dot_color, True, FONT_BODY)
        # bold lead before " : "
        if " — " in it:
            head, tail = it.split(" — ", 1)
            _set_run(p.add_run(), head + " — ", size, color, True, FONT_BODY)
            _set_run(p.add_run(), tail, size, color, False, FONT_BODY)
        else:
            _set_run(p.add_run(), it, size, color, False, FONT_BODY)
    return tb


def header(s, kicker, title, n):
    # accent bar
    rect(s, Inches(0.0), Inches(0.0), Inches(0.18), SH, fill=TEAL)
    txt(s, Inches(0.7), Inches(0.42), Inches(11.5), Inches(0.3),
        kicker.upper(), size=11.5, color=TEAL, bold=True, font=FONT_BODY)
    txt(s, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.9),
        title, size=27, color=INK, bold=True, font=FONT_TITLE, line_spacing=1.0)
    rect(s, Inches(0.72), Inches(1.62), Inches(0.9), Pt(3), fill=TEAL)
    # page number
    txt(s, Inches(12.5), Inches(7.02), Inches(0.7), Inches(0.3),
        str(n), size=10, color=MUTE, align=PP_ALIGN.RIGHT, font=FONT_MONO)
    txt(s, Inches(0.7), Inches(7.02), Inches(6), Inches(0.3),
        "Mnémo · Cas d'usage", size=9.5, color=MUTE, font=FONT_BODY)


def card(s, x, y, w, h, fill=WHITE, line=LINE, line_w=1.0):
    c = rect(s, x, y, w, h, fill=fill, line=line, line_w=line_w,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        c.adjustments[0] = 0.045
    except Exception:
        pass
    return c


def tag(s, x, y, w, text, fill, fg=WHITE, h=Inches(0.34), size=10.5):
    t = rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        t.adjustments[0] = 0.5
    except Exception:
        pass
    tf = t.text_frame; tf.word_wrap = False
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), text, size, fg, True, FONT_BODY)
    return t


# ===========================================================================
# SLIDE 1 — Couverture
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
# subtle accent block
rect(s, 0, Inches(6.55), SW, Inches(0.95), fill=TEAL_DARK)
rect(s, Inches(0.9), Inches(1.7), Inches(0.95), Pt(4), fill=TEAL)
txt(s, Inches(0.9), Inches(1.95), Inches(11), Inches(0.4),
    "CAS D'USAGE · COHORTE ENTREPRENEURS ÉLITE IA · LE LABO IA",
    size=13, color=TEAL_SOFT, bold=True, font=FONT_BODY)
txt(s, Inches(0.86), Inches(2.55), Inches(11.6), Inches(1.7),
    "Mnémo", size=66, color=WHITE, bold=True, font=FONT_TITLE, line_spacing=0.95)
txt(s, Inches(0.9), Inches(3.75), Inches(11.4), Inches(1.2),
    "Modèle économique & monétisation", size=30, color=WHITE, bold=True,
    font=FONT_TITLE, line_spacing=1.0)
txt(s, Inches(0.9), Inches(4.75), Inches(10.8), Inches(1.0),
    "Pourquoi nous avons fait évoluer le produit — les pivots de modèle "
    "économique et la correction de ses faiblesses.",
    size=16, color=RGBColor(0xC9, 0xD6, 0xD6), font=FONT_BODY, line_spacing=1.2)
txt(s, Inches(0.9), Inches(6.72), Inches(8), Inches(0.5),
    "La base mémorielle souveraine qui grandit avec votre organisation.",
    size=12.5, color=TEAL_SOFT, italic=True, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# SLIDE 2 — Le constat / le problème
# ===========================================================================
s = slide()
header(s, "Le constat", "La recommandation, seule, ne vaut plus rien", 2)
txt(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.7),
    "Demander « quelle infra de mémoire IA me faut-il ? » est devenu gratuit : "
    "un grand modèle de langage répond en quelques secondes. Conseiller ne suffit "
    "donc plus à construire une entreprise.",
    size=15.5, color=SLATE, line_spacing=1.25)

cols = [
    ("La reco est commoditisée",
     "N'importe quel LLM généraliste produit une recommandation crédible, gratuitement. "
     "Aucune barrière concurrentielle à ce niveau."),
    ("Les options existantes sont bancales",
     "Outils grand public (mémoire faible, données chez un tiers), projet sur-mesure "
     "(30–200 k€, non réutilisable) ou SaaS locataire à vie."),
    ("Le vrai problème est ailleurs",
     "Exécuter, maintenir dans le temps et ne pas se faire verrouiller par le "
     "fournisseur : c'est là que se cache la valeur défendable."),
]
cw = Inches(3.78); gapx = Inches(0.27); x0 = Inches(0.7); cy = Inches(3.0); chh = Inches(3.1)
for i, (h, b) in enumerate(cols):
    x = x0 + i * (cw + gapx)
    card(s, x, cy, cw, chh)
    rect(s, x + Inches(0.32), cy + Inches(0.35), Inches(0.55), Pt(4),
         fill=TEAL if i < 2 else GOLD)
    txt(s, x + Inches(0.32), cy + Inches(0.55), cw - Inches(0.64), Inches(1.0),
        h, size=16.5, color=INK, bold=True, font=FONT_TITLE, line_spacing=1.05)
    txt(s, x + Inches(0.32), cy + Inches(1.5), cw - Inches(0.64), Inches(1.5),
        b, size=13, color=SLATE, line_spacing=1.22)

txt(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
    "La question stratégique : que vend-on quand le conseil est gratuit ?",
    size=15, color=TEAL_DARK, bold=True, font=FONT_TITLE)

# ===========================================================================
# SLIDE 3 — Mnémo en une phrase
# ===========================================================================
s = slide()
header(s, "La réponse", "Recette ouverte, cuisine payante", 3)
card(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(1.55), fill=TEAL_SOFT, line=None)
txt(s, Inches(1.1), Inches(2.2), Inches(11.1), Inches(1.3),
    "Mnémo profile le besoin d'un dirigeant, recommande une stack de base mémorielle "
    "souveraine sur 7 couches — chiffrée et exportable — puis, en montée en gamme, "
    "la déploie et l'exploite via un agent.",
    size=16.5, color=INK, line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

duo = [
    ("La recette — OUVERTE & GRATUITE", TEAL,
     "On montre tout le « quoi » : la reco, la carte de coûts, le livrable exportable, "
     "le bundle de sortie reproductible. La transparence crée la confiance et tue le "
     "verrouillage.", "L'hameçon"),
    ("La cuisine — PAYANTE & RÉCURRENTE", TEAL_DARK,
     "On facture l'opération continue : déploiement assisté, monitoring, recalibration, "
     "négociation tarifaire au nom du client. C'est là que rentre l'argent récurrent.",
     "La rente"),
]
cw = Inches(5.82); x0 = Inches(0.7); cy = Inches(3.9); chh = Inches(2.7)
for i, (h, hc, b, badge) in enumerate(duo):
    x = x0 + i * (cw + Inches(0.26))
    card(s, x, cy, cw, chh)
    rect(s, x, cy, Inches(0.16), chh, fill=hc)
    txt(s, x + Inches(0.4), cy + Inches(0.32), cw - Inches(0.7), Inches(0.5),
        h, size=15, color=hc, bold=True, font=FONT_TITLE)
    txt(s, x + Inches(0.4), cy + Inches(1.0), cw - Inches(0.7), Inches(1.4),
        b, size=13.5, color=SLATE, line_spacing=1.25)
    tag(s, x + Inches(0.4), cy + chh - Inches(0.6), Inches(1.6), badge,
        fill=hc, size=11)

# ===========================================================================
# SLIDE 4 — Vue d'ensemble : ce qui a changé (synthèse des pivots)
# ===========================================================================
s = slide()
header(s, "Vue d'ensemble", "Ce qui a changé depuis la 1ʳᵉ version", 4)
txt(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.55),
    "Sept pivots ont transformé une simple idée de conseil en un produit défendable. "
    "Chacun corrige une faiblesse identifiée.",
    size=14.5, color=SLATE, line_spacing=1.2)

pivots = [
    ("01", "Version mince refusée", "Production-ready par lots, zéro dette — crédibilité."),
    ("02", "On ne vend plus la reco", "« Recette ouverte, cuisine payante » : on facture l'opération."),
    ("03", "Mandataire fiduciaire", "Payé par le client, jamais commissionné par les fournisseurs."),
    ("04", "Trio d'avantages défendables", "Sortie certifiée + mode fiduciaire + effet réseau de données."),
    ("05", "Honnêteté ±30 % assumée", "Slider de projection + prix dynamiques, pas de fausse précision."),
    ("06", "Agent supervisé par design", "Jamais de création de compte ni de carte — risque juridique levé."),
    ("07", "Effet réseau « en rails »", "La donnée se densifie dès les premiers déploiements monitorés."),
]
cw = Inches(3.78); gapx = Inches(0.27); x0 = Inches(0.7)
ys = [Inches(2.75), Inches(4.0), Inches(5.25)]
chh = Inches(1.12)
for idx, (n, h, b) in enumerate(pivots):
    row = idx // 3; col = idx % 3
    x = x0 + col * (cw + gapx); y = ys[row]
    card(s, x, y, cw, chh)
    txt(s, x + Inches(0.28), y + Inches(0.16), Inches(0.9), Inches(0.5),
        n, size=20, color=TEAL, bold=True, font=FONT_MONO)
    txt(s, x + Inches(1.0), y + Inches(0.16), cw - Inches(1.25), Inches(0.4),
        h, size=13.5, color=INK, bold=True, font=FONT_TITLE)
    txt(s, x + Inches(1.0), y + Inches(0.55), cw - Inches(1.25), Inches(0.5),
        b, size=10.8, color=SLATE, line_spacing=1.12)
# 8th cell : leçon
x = x0 + 1 * (cw + gapx); y = ys[2]  # placeholder not used
# place a wide takeaway under grid
txt(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5),
    "Fil conducteur : passer d'une idée copiable à un système où la confiance, "
    "l'anti-verrouillage et la donnée deviennent des barrières concurrentielles.",
    size=13.5, color=TEAL_DARK, bold=True, font=FONT_TITLE, line_spacing=1.15)

# ===========================================================================
# Helper — slide « avant / après »
# ===========================================================================
def pivot_slide(n, kicker, title, rows, lesson):
    s = slide()
    header(s, kicker, title, n)
    # column headers
    top = Inches(2.0)
    lw = Inches(5.82); x0 = Inches(0.7)
    xr = x0 + lw + Inches(0.26)
    tag(s, x0, top, Inches(1.7), "AVANT", fill=RED_SOFT, size=11)
    tag(s, xr, top, Inches(1.7), "APRÈS", fill=TEAL, size=11)
    ry = Inches(2.55)
    rh = Inches(1.12)
    for (av, ap) in rows:
        card(s, x0, ry, lw, rh, fill=WHITE, line=LINE)
        card(s, xr, ry, lw, rh, fill=TEAL_SOFT, line=None)
        txt(s, x0 + Inches(0.3), ry + Inches(0.18), lw - Inches(0.6), rh - Inches(0.3),
            av, size=13, color=SLATE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
        txt(s, xr + Inches(0.3), ry + Inches(0.18), lw - Inches(0.6), rh - Inches(0.3),
            ap, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18, bold=False)
        ry = ry + rh + Inches(0.18)
    # lesson band
    by = ry + Inches(0.05)
    if by > Inches(6.35):
        by = Inches(6.35)
    card(s, Inches(0.7), by, Inches(11.9), Inches(0.78), fill=INK, line=None)
    txt(s, Inches(1.05), by, Inches(11.3), Inches(0.78),
        lesson, size=13.5, color=WHITE, bold=True, font=FONT_TITLE,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    return s

# ===========================================================================
# SLIDE 5 — Pivot A : modèle de vente
# ===========================================================================
pivot_slide(
    5, "Pivot · modèle économique",
    "De « vendre la reco » à « cuisine payante »",
    [
        ("La valeur = la recommandation, vendue une fois.",
         "La valeur = l'opération continue, facturée en récurrent."),
        ("Faiblesse : un LLM produit la même reco gratuitement.",
         "On ouvre la recette (transparence) et on facture la cuisine."),
        ("Le livrable était notre produit — donc copiable.",
         "Le livrable devient l'hameçon ; la rente est ailleurs."),
    ],
    "Faiblesse corrigée : la commoditisation.  Impact : un revenu récurrent, pas un one-shot."
)

# ===========================================================================
# SLIDE 6 — Pivot B : fiduciaire
# ===========================================================================
pivot_slide(
    6, "Pivot · monétisation",
    "Le mandataire fiduciaire, condition de survie",
    [
        ("Tentation : se faire commissionner par les fournisseurs.",
         "Payé uniquement par le client, jamais par les fournisseurs."),
        ("Faiblesse : conflit d'intérêt, perte de confiance.",
         "Rémunération affichée ; négociation tarifaire au nom du client."),
        ("Modèle des comparateurs « gratuits » qui orientent vers qui paie.",
         "Engagement contractuel : zéro commission cachée."),
    ],
    "Faiblesse corrigée : le conflit d'intérêt.  Un mandataire payé par le fournisseur trahit le mandat — et meurt."
)

# ===========================================================================
# SLIDE 7 — Pivot C : honnêteté du coût + agent supervisé
# ===========================================================================
pivot_slide(
    7, "Pivot · produit & risque",
    "Honnêteté des coûts & agent supervisé",
    [
        ("Prétendre prédire le coût exact (fausse précision).",
         "Honnêteté ±30 % assumée : slider de projection + fourchettes."),
        ("Table de prix figée → dette de maintenance permanente.",
         "Flux de prix dynamique : la dette devient un flux automatisé."),
        ("Agent autonome qui crée comptes et saisit la carte.",
         "Agent supervisé : jamais de compte ni de carte — risque juridique levé."),
    ],
    "Faiblesses corrigées : fausse rigueur, dette fournisseur et risque de responsabilité."
)

# ===========================================================================
# SLIDE 8 — Modèle économique : gratuit vs payant (tableau)
# ===========================================================================
s = slide()
header(s, "Modèle économique", "Gratuit (l'hameçon) vs payant (la rente)", 8)

lw = Inches(5.82); x0 = Inches(0.7); xr = x0 + lw + Inches(0.26)
top = Inches(2.0)
# headers
hL = rect(s, x0, top, lw, Inches(0.62), fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
hL.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
txt(s, x0 + Inches(0.35), top, lw - Inches(0.6), Inches(0.62),
    "GRATUIT — créer l'habitude", size=15, color=WHITE, bold=True,
    font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)
hR = rect(s, xr, top, lw, Inches(0.62), fill=TEAL_DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, xr + Inches(0.35), top, lw - Inches(0.6), Inches(0.62),
    "PAYANT — le revenu récurrent", size=15, color=WHITE, bold=True,
    font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

free = [
    "Recommandation de stack sur 7 couches",
    "Carte de coûts projetée (fourchettes sourcées)",
    "Livrable exportable (plan + stack)",
    "Bundle de sortie reproductible (anti-verrouillage)",
]
paid = [
    "Déploiement assisté de l'infrastructure",
    "Monitoring continu & recalibration",
    "Négociation tarifaire au nom du client",
    "Bundle de sortie premium & garanties",
]
cy = top + Inches(0.78)
card(s, x0, cy, lw, Inches(3.55), fill=WHITE)
card(s, xr, cy, lw, Inches(3.55), fill=TEAL_SOFT, line=None)
bullets(s, x0 + Inches(0.35), cy + Inches(0.3), lw - Inches(0.7), Inches(3.0),
        free, size=14, gap=13, dot_color=TEAL)
bullets(s, xr + Inches(0.35), cy + Inches(0.3), lw - Inches(0.7), Inches(3.0),
        paid, size=14, gap=13, dot_color=TEAL_DARK)
# footers
txt(s, x0 + Inches(0.35), cy + Inches(3.05), lw - Inches(0.7), Inches(0.4),
    "Aucune rétention d'information : tout le « quoi » est offert.",
    size=11.5, color=MUTE, italic=True)
txt(s, xr + Inches(0.35), cy + Inches(3.05), lw - Inches(0.7), Inches(0.4),
    "On facture le « faire » et le « maintenir », pas le « savoir ».",
    size=11.5, color=TEAL_DARK, italic=True, bold=True)

txt(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.5),
    "Logique du palier gratuit généreux : créer l'habitude, puis monétiser l'exploitation.",
    size=13, color=SLATE, font=FONT_TITLE, bold=True)

# ===========================================================================
# SLIDE 9 — Monétisation : flux de revenus & rôle fiduciaire
# ===========================================================================
s = slide()
header(s, "Monétisation", "Flux de revenus & rôle fiduciaire", 9)

cols = [
    ("Ce qui est ponctuel", TEAL, [
        "Bundle de sortie premium",
        "Mise en place initiale du déploiement",
    ], "Déclencheur d'entrée"),
    ("Ce qui est récurrent", TEAL_DARK, [
        "Abonnement de monitoring",
        "Recalibration continue de la stack",
        "Re-optimisation sur économie validée",
    ], "Le cœur de la rente"),
    ("Le rôle fiduciaire", GOLD, [
        "Négociation tarifaire au nom du client",
        "Rémunération affichée, zéro commission cachée",
        "Aligne nos revenus sur ceux du client",
    ], "La confiance monétisée"),
]
cw = Inches(3.78); gapx = Inches(0.27); x0 = Inches(0.7); cy = Inches(2.1); chh = Inches(3.7)
for i, (h, hc, items, badge) in enumerate(cols):
    x = x0 + i * (cw + gapx)
    card(s, x, cy, cw, chh)
    rect(s, x, cy, cw, Inches(0.12), fill=hc)
    txt(s, x + Inches(0.3), cy + Inches(0.35), cw - Inches(0.6), Inches(0.5),
        h, size=16, color=hc, bold=True, font=FONT_TITLE)
    bullets(s, x + Inches(0.3), cy + Inches(1.05), cw - Inches(0.6), Inches(2.1),
            items, size=12.5, gap=9, dot_color=hc)
    tag(s, x + Inches(0.3), cy + chh - Inches(0.6), Inches(2.7), badge,
        fill=hc, size=10.5)

txt(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.85),
    "Le récurrent est la colonne vertébrale du modèle : la valeur se mesure dans la "
    "durée, pas à la signature. Le rôle fiduciaire transforme la confiance en ligne de revenu.",
    size=13.5, color=TEAL_DARK, bold=True, font=FONT_TITLE, line_spacing=1.2)

# ===========================================================================
# SLIDE 10 — Le flywheel (effet réseau)
# ===========================================================================
s = slide()
header(s, "Effet réseau", "Le flywheel de la donnée de coût réel", 10)
txt(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.55),
    "Chaque déploiement monitoré nourrit le suivant. Une donnée que personne ne peut répliquer.",
    size=14.5, color=SLATE, line_spacing=1.2)

steps = [
    ("Déploiement\nmonitoré", TEAL),
    ("Coût réel\nanonymisé remonté", BLUE),
    ("La reco du client\nsuivant s'affine", TEAL),
    ("Donnée\nirréplicable", GOLD),
    ("Plus de clients\n→ retour à l'étape 1", TEAL_DARK),
]
# circular-ish horizontal flow
bw = Inches(2.05); bh = Inches(1.3); cy = Inches(3.1)
gap = Inches(0.32)
x = Inches(0.72)
for i, (label, c) in enumerate(steps):
    b = card(s, x, cy, bw, bh, fill=WHITE, line=c, line_w=1.6)
    txt(s, x + Inches(0.12), cy, bw - Inches(0.24), bh,
        label, size=13, color=INK, bold=True, font=FONT_TITLE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    # number badge
    nb = rect(s, x + Inches(0.08), cy - Inches(0.18), Inches(0.36), Inches(0.36),
              fill=c, shape=MSO_SHAPE.OVAL)
    txt(s, x + Inches(0.08), cy - Inches(0.18), Inches(0.36), Inches(0.36),
        str(i + 1), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    if i < len(steps) - 1:
        ar = rect(s, x + bw + Inches(0.02), cy + bh/2 - Inches(0.14),
                  gap - Inches(0.04), Inches(0.28), fill=TEAL,
                  shape=MSO_SHAPE.RIGHT_ARROW)
    x = x + bw + gap

# loop-back arrow band
rect(s, Inches(0.72), Inches(4.75), Inches(11.85), Inches(0.06), fill=LINE)
card(s, Inches(0.72), Inches(5.05), Inches(11.85), Inches(1.45), fill=TEAL_SOFT, line=None)
txt(s, Inches(1.1), Inches(5.25), Inches(11.1), Inches(1.1),
    "Plus le réseau grossit, plus les recommandations sont précises — et plus la "
    "précision attire de clients. Cette boucle de données est la barrière concurrentielle "
    "la plus profonde : elle se construit dans le temps et ne se copie pas.",
    size=14.5, color=INK, line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# SLIDE 11 — Les 3 avantages défendables
# ===========================================================================
s = slide()
header(s, "Avantages défendables", "Trois remparts, pas une simple reco", 11)
txt(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.5),
    "La recommandation seule n'est pas défendable. Trois mécanismes le rendent — chacun "
    "inspiré d'une autre industrie.",
    size=14.5, color=SLATE, line_spacing=1.2)

advs = [
    ("Sortie certifiée", TEAL,
     "En un clic, un bundle reproductible redéployable ailleurs.",
     "Défend contre la capture — l'anti-verrouillage devient contractuel."),
    ("Mode fiduciaire", TEAL_DARK,
     "Payé par le client, jamais commissionné par les fournisseurs.",
     "Défend la confiance — un concurrent commissionné ne peut pas copier."),
    ("Effet réseau de données", GOLD,
     "Le coût réel anonymisé remonté par chaque déploiement.",
     "Défend par la donnée — une calibration que nul ne peut répliquer."),
]
cw = Inches(3.78); gapx = Inches(0.27); x0 = Inches(0.7); cy = Inches(2.75); chh = Inches(3.55)
for i, (h, c, what, why) in enumerate(advs):
    x = x0 + i * (cw + gapx)
    card(s, x, cy, cw, chh)
    circ = rect(s, x + Inches(0.3), cy + Inches(0.35), Inches(0.62), Inches(0.62),
                fill=c, shape=MSO_SHAPE.OVAL)
    txt(s, x + Inches(0.3), cy + Inches(0.35), Inches(0.62), Inches(0.62),
        str(i + 1), size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    txt(s, x + Inches(0.3), cy + Inches(1.15), cw - Inches(0.6), Inches(0.6),
        h, size=17, color=INK, bold=True, font=FONT_TITLE, line_spacing=1.0)
    txt(s, x + Inches(0.3), cy + Inches(1.8), cw - Inches(0.6), Inches(1.0),
        what, size=12.8, color=SLATE, line_spacing=1.22)
    rect(s, x + Inches(0.3), cy + chh - Inches(1.15), cw - Inches(0.6), Pt(2), fill=LINE)
    txt(s, x + Inches(0.3), cy + chh - Inches(1.0), cw - Inches(0.6), Inches(0.85),
        why, size=12, color=c, bold=True, line_spacing=1.18, font=FONT_BODY)

# ===========================================================================
# SLIDE 12 — Cibles & mise sur le marché
# ===========================================================================
s = slide()
header(s, "Cibles & mise sur le marché", "À qui l'on vend, et dans quel ordre", 12)

stages = [
    ("1.", "Bâtisseurs souverains", TEAL,
     "La communauté d'Amine : dirigeants, coachs et consultants qui veulent leur base "
     "mémorielle souveraine — revendable. Accès direct, demande déjà acquise, bouche-à-oreille.",
     "Tête de pont"),
    ("2.", "PME & startups tech", BLUE,
     "Équipes de 5 à 10 personnes, technophiles, cycle de vente court. Le segment "
     "naturel d'expansion une fois la preuve faite.",
     "Expansion"),
    ("3.", "Cabinets régulés & programmes", GOLD,
     "Conformité forte = argument premium, mais cycle long. Volume et multi-organisations "
     "viennent plus tard. Grands groupes : hors périmètre initial.",
     "Plus tard"),
]
cy = Inches(2.1); rh = Inches(1.42)
for i, (n, h, c, b, badge) in enumerate(stages):
    y = cy + i * (rh + Inches(0.16))
    card(s, Inches(0.7), y, Inches(11.9), rh)
    rect(s, Inches(0.7), y, Inches(0.16), rh, fill=c)
    txt(s, Inches(1.05), y + Inches(0.2), Inches(0.7), Inches(1.0),
        n, size=26, color=c, bold=True, font=FONT_MONO)
    txt(s, Inches(1.85), y + Inches(0.22), Inches(3.3), Inches(1.0),
        h, size=16.5, color=INK, bold=True, font=FONT_TITLE, line_spacing=1.05,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.25), y + Inches(0.18), Inches(5.6), rh - Inches(0.3),
        b, size=12.8, color=SLATE, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
    tag(s, Inches(10.95), y + rh/2 - Inches(0.17), Inches(1.45), badge, fill=c, size=10)

# ===========================================================================
# SLIDE 13 — Risques neutralisés par conception
# ===========================================================================
s = slide()
header(s, "Robustesse", "Risques neutralisés par conception", 13)
txt(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.5),
    "Chaque faiblesse du produit a une réponse intégrée dès la conception — pas un correctif tardif.",
    size=14.5, color=SLATE, line_spacing=1.2)

risks = [
    ("Conflit d'intérêt fournisseur", "Mode fiduciaire : zéro commission cachée, divulgation contractuelle."),
    ("Le livrable sabote la vente", "Recette ouverte assumée ; la rente, c'est l'opération continue."),
    ("Fausse précision des coûts", "Fourchettes ±30 % + sources datées + recalibration par le réseau."),
    ("Dette de maintenance fournisseur", "Flux de prix automatisé ; jamais de table figée."),
    ("Risque juridique de l'agent", "L'agent ne crée jamais de compte ni ne saisit de carte."),
    ("Démarrage à froid de la donnée", "Effet réseau « en rails » : se densifie dès le 1ᵉʳ déploiement."),
]
cw = Inches(5.82); gapx = Inches(0.26); x0 = Inches(0.7)
cy = Inches(2.75); rh = Inches(1.18)
for i, (r, m) in enumerate(risks):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gapx); y = cy + row * (rh + Inches(0.16))
    card(s, x, y, cw, rh)
    # check chip
    chip = rect(s, x + Inches(0.28), y + Inches(0.3), Inches(0.4), Inches(0.4),
                fill=TEAL_SOFT, shape=MSO_SHAPE.OVAL)
    txt(s, x + Inches(0.28), y + Inches(0.28), Inches(0.4), Inches(0.4),
        "✓", size=15, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.85), y + Inches(0.18), cw - Inches(1.1), Inches(0.4),
        r, size=13.5, color=INK, bold=True, font=FONT_TITLE)
    txt(s, x + Inches(0.85), y + Inches(0.58), cw - Inches(1.1), Inches(0.55),
        m, size=11.8, color=SLATE, line_spacing=1.15)

# ===========================================================================
# SLIDE 14 — Roadmap : quand l'argent rentre
# ===========================================================================
s = slide()
header(s, "Roadmap", "Trois lots — et quand l'argent rentre", 14)

lots = [
    ("LOT 1", "Conseil + remparts", TEAL,
     ["Profilage, reco, carte de coûts", "Livrable & bundle de sortie",
      "Mode fiduciaire & sortie certifiée"],
     "Vendable seul. Crée l'habitude.", "Argent : ponctuel"),
    ("LOT 2", "La cuisine payante", TEAL_DARK,
     ["Déploiement assisté", "Monitoring continu & recalibration",
      "Effet réseau activé"],
     "Le revenu récurrent rentre. Le flywheel s'allume.", "Argent : RÉCURRENT"),
    ("LOT 3", "L'écosystème", GOLD,
     ["Migration garantie", "Re-optimisation continue", "Produits dérivés"],
     "Expansion de la valeur dans la durée.", "Argent : expansion"),
]
cw = Inches(3.78); gapx = Inches(0.27); x0 = Inches(0.7); cy = Inches(2.05); chh = Inches(4.05)
for i, (lot, h, c, items, note, money) in enumerate(lots):
    x = x0 + i * (cw + gapx)
    card(s, x, cy, cw, chh)
    rect(s, x, cy, cw, Inches(0.78), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, cy + Inches(0.4), cw, Inches(0.38), fill=c)  # square off bottom of header
    txt(s, x + Inches(0.3), cy + Inches(0.1), cw - Inches(0.6), Inches(0.3),
        lot, size=12, color=WHITE, bold=True, font=FONT_MONO)
    txt(s, x + Inches(0.3), cy + Inches(0.36), cw - Inches(0.6), Inches(0.42),
        h, size=16, color=WHITE, bold=True, font=FONT_TITLE)
    bullets(s, x + Inches(0.3), cy + Inches(1.05), cw - Inches(0.6), Inches(1.7),
            items, size=12.5, gap=9, dot_color=c)
    rect(s, x + Inches(0.3), cy + Inches(2.85), cw - Inches(0.6), Pt(2), fill=LINE)
    txt(s, x + Inches(0.3), cy + Inches(3.0), cw - Inches(0.6), Inches(0.55),
        note, size=11.8, color=SLATE, italic=True, line_spacing=1.15)
    tag(s, x + Inches(0.3), cy + chh - Inches(0.58), Inches(2.6), money, fill=c, size=10.5)

txt(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5),
    "Le Lot 2 est le moment charnière : c'est là que le récurrent démarre et que l'effet réseau s'enclenche.",
    size=13, color=TEAL_DARK, bold=True, font=FONT_TITLE)

# ===========================================================================
# SLIDE 15 — Synthèse / la leçon pour la cohorte
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, Inches(0.18), SH, fill=TEAL)
txt(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.4),
    "SYNTHÈSE · LA LEÇON POUR LA COHORTE", size=12.5, color=TEAL_SOFT, bold=True)
txt(s, Inches(0.86), Inches(1.15), Inches(11.6), Inches(1.0),
    "Ne vendez pas ce qu'un LLM offre gratuitement.", size=30, color=WHITE,
    bold=True, font=FONT_TITLE, line_spacing=1.02)

takeaways = [
    ("Donnez la recette, facturez la cuisine.",
     "Le conseil est commoditisé. La valeur durable est dans l'exécution et le maintien."),
    ("La confiance se monétise.",
     "Le mode fiduciaire — payé par le client, jamais par les fournisseurs — est une condition de survie, pas une option."),
    ("Construisez une barrière qui grossit toute seule.",
     "L'effet réseau de données et l'anti-verrouillage se renforcent avec chaque client."),
    ("Corrigez les faiblesses par conception.",
     "Honnêteté des coûts, agent supervisé, risques neutralisés : la robustesse est un argument de vente."),
]
cy = Inches(2.5); rh = Inches(1.0)
for i, (h, b) in enumerate(takeaways):
    y = cy + i * (rh + Inches(0.12))
    rect(s, Inches(0.9), y + Inches(0.08), Inches(0.5), Inches(0.5),
         fill=TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.9), y + Inches(0.06), Inches(0.5), Inches(0.5),
        str(i + 1), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
    txt(s, Inches(1.6), y, Inches(10.8), Inches(0.42),
        h, size=16.5, color=WHITE, bold=True, font=FONT_TITLE)
    txt(s, Inches(1.6), y + Inches(0.42), Inches(10.8), Inches(0.55),
        b, size=12.8, color=RGBColor(0xC9, 0xD6, 0xD6), line_spacing=1.15)

rect(s, 0, Inches(7.0), SW, Inches(0.5), fill=TEAL_DARK)
txt(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.5),
    "Mnémo — la base mémorielle souveraine qui grandit avec votre organisation.",
    size=12.5, color=TEAL_SOFT, italic=True, anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------------------
out = r"C:\Users\amans\OneDrive\Projets\Infra\presentation\Mnemo-modele-economique.pptx"
prs.save(out)
print("OK", out, len(prs.slides.__iter__.__self__._sldIdLst))
