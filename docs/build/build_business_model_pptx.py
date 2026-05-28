# -*- coding: utf-8 -*-
"""
Génère le deck « Business model — Strate » (docs/strate-business-model.pptx).

Design épuré aligné sur design-reference/mn_mo_brand_identity/DESIGN.md :
  - Primaire teal #00685F, secondaire bleu #0051D5, premium or #7D5400 / #9D6A00
  - Surface #FAF8FF, texte #131B2E, variante #3D4947, contour #BCC9C6
  - Titres = Space Grotesk, corps = Inter, données/chiffres = JetBrains Mono

Principe : la slide reste épurée (titre court + 3-5 puces) ; le détail, les
chiffres et les arguments vivent dans les NOTES du présentateur.

Toutes les prévisions chiffrées sont étiquetées HYPOTHÈSE. Le prix de vente
Strate est un [PLACEHOLDER] (sondage Van Westendorp + conjoint en cours,
cf. docs/pricing/wtp-research.md). Aucune stat de marché n'est inventée.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Palette (DESIGN.md)
# ---------------------------------------------------------------------------
TEAL        = RGBColor(0x00, 0x68, 0x5F)   # primary
TEAL_DIM    = RGBColor(0x00, 0x83, 0x78)   # primary-container
BLUE        = RGBColor(0x00, 0x51, 0xD5)   # secondary
GOLD        = RGBColor(0x9D, 0x6A, 0x00)   # tertiary-container (premium)
GOLD_TXT    = RGBColor(0x7D, 0x54, 0x00)   # tertiary
INK         = RGBColor(0x13, 0x1B, 0x2E)   # on-surface
INK_VARIANT = RGBColor(0x3D, 0x49, 0x47)   # on-surface-variant
SURFACE     = RGBColor(0xFA, 0xF8, 0xFF)   # surface
SURFACE_LOW = RGBColor(0xF2, 0xF3, 0xFF)   # surface-container-low
CARD_BORDER = RGBColor(0xBC, 0xC9, 0xC6)   # outline-variant
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ERROR       = RGBColor(0xBA, 0x1A, 0x1A)   # error
INVERSE     = RGBColor(0x28, 0x30, 0x44)   # inverse-surface (dark band)
INVERSE_TXT = RGBColor(0xEE, 0xF0, 0xFF)   # inverse-on-surface

FONT_TITLE = "Space Grotesk"
FONT_BODY  = "Inter"
FONT_MONO  = "JetBrains Mono"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _no_autofit(tf):
    # désactive l'autosize pour garder une mise en page maîtrisée
    tf.word_wrap = True


def add_slide(bg=SURFACE):
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, bg)
    return s


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size, color, font=FONT_BODY, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def para(tf, text, size, color, font=FONT_BODY, bold=False, italic=False,
         space_after=6, space_before=0, align=PP_ALIGN.LEFT, level=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    set_run(r, text, size, color, font, bold, italic)
    return p


def bullet(tf, text, size=15, color=INK_VARIANT, bold_lead=None,
           space_after=9, first=False, marker="—", marker_color=TEAL):
    """Puce épurée : tiret coloré + texte. bold_lead = segment en gras au début."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = 1.08
    rm = p.add_run()
    set_run(rm, marker + "  ", size, marker_color, FONT_BODY, bold=True)
    if bold_lead:
        rb = p.add_run()
        set_run(rb, bold_lead, size, INK, FONT_BODY, bold=True)
    rt = p.add_run()
    set_run(rt, text, size, color, FONT_BODY)
    return p


def accent_bar(slide, color=TEAL, t=Inches(0.0)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), t, Inches(0.16), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def eyebrow(slide, text, color=TEAL, l=Inches(0.85), t=Inches(0.55), w=Inches(11.6)):
    tb, tf = textbox(slide, l, t, w, Inches(0.35))
    p = para(tf, text.upper(), 12.5, color, FONT_BODY, bold=True, first=True)
    # letter spacing
    for r in p.runs:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', '140')
    return tb


def title(slide, text, l=Inches(0.85), t=Inches(0.92), w=Inches(11.6), h=Inches(1.0), size=30, color=INK):
    tb, tf = textbox(slide, l, t, w, h)
    para(tf, text, size, color, FONT_TITLE, bold=True, first=True)
    return tb


def rule(slide, l=Inches(0.85), t=Inches(1.78), w=Inches(11.6), color=CARD_BORDER):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(1.4))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def card(slide, l, t, w, h, fill=WHITE, border=CARD_BORDER, border_w=1.0, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    c = slide.shapes.add_shape(shape_type, l, t, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    if border is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = border
        c.line.width = Pt(border_w)
    c.shadow.inherit = False
    # adoucir le rayon
    if radius:
        try:
            c.adjustments[0] = 0.06
        except Exception:
            pass
    return c


def card_text(card_shape, blocks, anchor=MSO_ANCHOR.TOP):
    """blocks = liste de dict {t, size, color, font, bold, italic, sa, sb, align}"""
    tf = card_shape.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.14)
    for i, b in enumerate(blocks):
        para(tf, b["t"], b.get("size", 14), b.get("color", INK_VARIANT),
             b.get("font", FONT_BODY), b.get("bold", False), b.get("italic", False),
             b.get("sa", 5), b.get("sb", 0), b.get("align", PP_ALIGN.LEFT),
             b.get("level", 0), first=(i == 0))


def pagefoot(slide, n, label="Strate · Business model"):
    tb, tf = textbox(slide, Inches(0.85), Inches(7.05), Inches(11.6), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, label, 9, CARD_BORDER, FONT_MONO)
    rn = p.add_run()
    set_run(rn, "", 9, CARD_BORDER, FONT_MONO)
    # numéro à droite
    tb2, tf2 = textbox(slide, Inches(11.8), Inches(7.05), Inches(0.65), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    set_run(r2, str(n).zfill(2), 9, CARD_BORDER, FONT_MONO)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ===========================================================================
# Construction du deck
# ===========================================================================
PAGE = 0


def std_header(slide, eb, ttl, eb_color=TEAL, accent=TEAL):
    accent_bar(slide, accent)
    eyebrow(slide, eb, eb_color)
    title(slide, ttl)
    rule(slide)


# --- 1. COUVERTURE ---------------------------------------------------------
PAGE += 1
s = add_slide(INVERSE)
accent_bar(s, TEAL)
# pastille / eyebrow clair
tb, tf = textbox(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.4))
p = para(tf, "BUSINESS MODEL · DOCUMENT DE TRAVAIL · 2026", 13, RGBColor(0x6B, 0xD8, 0xCB), FONT_BODY, bold=True, first=True)
for r in p.runs:
    r._r.get_or_add_rPr().set('spc', '180')
tb, tf = textbox(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.6))
para(tf, "Strate", 66, WHITE, FONT_TITLE, bold=True, first=True, space_after=2)
para(tf, "La base mémorielle souveraine qui grandit avec votre organisation —", 21, INVERSE_TXT, FONT_TITLE, bold=False, space_after=0)
para(tf, "sans migration, sans verrouillage, sans coûts cachés.", 21, INVERSE_TXT, FONT_TITLE, bold=False)
tb, tf = textbox(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(1.4))
para(tf, "Configurateur SaaS souverain d'infrastructure de « base mémorielle » IA.", 15, INVERSE_TXT, FONT_BODY, first=True, space_after=4)
para(tf, "Conseil → déploiement · recette ouverte, cuisine payante · trio de moats : Exit Escrow · Fiduciary Mode · Intelligence Network.", 14, RGBColor(0xB4, 0xC5, 0xFF), FONT_BODY)
tb, tf = textbox(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))
para(tf, "Prix de vente = [PLACEHOLDER] · sondage Van Westendorp + conjoint en cours · toute projection chiffrée = HYPOTHÈSE à valider.", 11, RGBColor(0x6B, 0xD8, 0xCB), FONT_MONO, first=True)
notes(s, """
Slide d'ouverture. Cadrer en une phrase : Strate est un configurateur SaaS souverain qui aide un dirigeant non-expert à choisir, chiffrer, exporter et (en upsell) déployer puis exploiter une « base mémorielle » IA souveraine, sans se faire capturer par les fournisseurs.

Ce que ce deck est : un business model de travail, profond, mais EXPLICITEMENT hypothétique sur tous les chiffres. Ce qu'il n'est pas : une promesse de revenus.

Trois avertissements à poser d'entrée (règle DÉFCON 1 d'Amine, zéro donnée erronée présentée comme fiable) :
1. Le prix de vente du service Strate est un [PLACEHOLDER] : un sondage de willingness-to-pay (Van Westendorp + conjoint CBC, cf. docs/pricing/wtp-research.md) est en cours. Aucun montant ferme ne sera asséné ici, seulement des fourchettes d'hypothèses étiquetées.
2. Tout P&L / toute taille de marché présenté plus loin est une hypothèse de modélisation, pas une certitude.
3. Le produit décrit est le produit RÉEL du dépôt (Lot 1 livré et déployé en prod, épic LLM complet) ; ce qui est « à venir » est signalé comme tel.

Le nom « Strate » : une strate = mémoire sédimentée qui se conserve dans le temps + une couche (l'infra = 7 couches C0→C6). Designer l'infra = empiler les bonnes strates.
""")
pagefoot(s, PAGE, "")

# --- 2. EXECUTIVE SUMMARY ---------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Synthèse", "La thèse, en une page")
tb, tf = textbox(s, Inches(0.85), Inches(2.05), Inches(11.6), Inches(0.95))
para(tf, "« Le Cloudflare de la mémoire IA souveraine : se mettre devant l'infra sans capturer le client, et faire de l'anti-lock-in un produit qui se vend. »", 19, TEAL, FONT_TITLE, bold=True, italic=True, first=True)
# 4 cartes points-clés
cw, ch, gap, top = Inches(2.78), Inches(2.55), Inches(0.18), Inches(3.35)
left0 = Inches(0.85)
pts = [
    ("Le problème", "Un dirigeant qui veut une mémoire IA n'a que des options bancales : grand public (données chez un tiers US), sur-mesure (30-200 k€), ou SaaS US (locataire à vie).", TEAL),
    ("La solution", "Profilage → reco 7 couches sourcée ±30 % → livrable → bundle de redéploiement (Exit Escrow). En upsell : déploiement assisté + supervision.", BLUE),
    ("Le moat", "Un trio défendable que les concurrents commissionnés ne peuvent pas copier : Exit Escrow, Fiduciary Mode, Intelligence Network.", GOLD_TXT),
    ("Le modèle", "Recette ouverte (gratuite), cuisine payante. Le service a un prix ferme ; l'infra reste au coût fournisseur, zéro marge cachée.", TEAL),
]
for i, (h, body, col) in enumerate(pts):
    l = Emu(int(left0) + i * (int(cw) + int(gap)))
    c = card(s, l, top, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, top, cw, Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    card_text(c, [
        {"t": h, "size": 16, "color": INK, "font": FONT_TITLE, "bold": True, "sa": 8, "sb": 6},
        {"t": body, "size": 12.5, "color": INK_VARIANT, "sa": 4},
    ])
pagefoot(s, PAGE)
notes(s, """
La thèse d'investissement en une phrase : reproduire le playbook Cloudflare sur l'infra de mémoire IA. Cloudflare s'est mis DEVANT l'origine sans capturer le client (on peut partir), a créé l'habitude par un free tier + transparence, et a transformé sa visibilité du trafic en meilleure défense pour tous (effet réseau). Strate réplique ces trois ressorts (Exit Escrow / Fiduciary / Network).

Points clés à marteler :
- Le problème est réel et structurel : perte de connaissance critique quand un sachant part ; re-contextualisation permanente dans chaque chat ; non-souveraineté des données ; vendor lock-in. (PRD §1, genèse P4.)
- La différenciation n'est PAS le moteur de reco (un LLM généraliste le commoditise) — c'est le trio de moats.
- Le modèle économique « recette ouverte, cuisine payante » désamorce le piège classique « le livrable gratuit sabote l'upsell » : on ne facture pas l'information (le QUOI), on facture l'opération continue (le déploiement, la supervision, la recalibration, la négociation vendor).

Honnêteté DÉFCON 1 : à ce stade, la PREUVE dure est l'anti-lock-in (Exit Escrow testable). Les autres bénéfices (gain de temps, qualité de réponse) sont « à mesurer chez vous », jamais chiffrés sans donnée.

Statut produit : Lot 1 (conseil + moats) livré, déployé en prod (infra.ai-mpower.com), épic LLM complet (intake libre, narration, assistant Q&A, reco vivante sourcée, console admin). Lot 2 (déploiement assisté + monitoring) et résidence/DR = à venir.
""")

# --- 3. PROBLÈME & JTBD -----------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Problème & Jobs-To-Be-Done", "On vend un résultat souverain, pas une stack", accent=TEAL)
# colonne gauche : les 4 jobs
tb, tf = textbox(s, Inches(0.85), Inches(2.05), Inches(6.05), Inches(4.6))
para(tf, "Les 4 jobs du dirigeant", 16, INK, FONT_TITLE, bold=True, first=True, space_after=10)
for j in [
    "Ne plus perdre la connaissance critique quand un sachant part.",
    "Que mon IA réponde avec MA doctrine, pas la moyenne du web.",
    "Garder la souveraineté sans payer un projet à 200 k€.",
    "Choisir la bonne infra sans être expert — et pouvoir partir quand je veux.",
]:
    bullet(tf, j, size=15, space_after=12)
# colonne droite : pourquoi maintenant
c = card(s, Inches(7.15), Inches(2.05), Inches(5.3), Inches(4.55), fill=SURFACE_LOW)
card_text(c, [
    {"t": "Pourquoi maintenant", "size": 16, "color": TEAL, "font": FONT_TITLE, "bold": True, "sa": 9, "sb": 2},
    {"t": "Souveraineté IA : exigence montante (UE, secteur public, régulés).", "size": 13, "sa": 7},
    {"t": "RGPD + CNDP (Maroc) : la donnée ne peut plus dormir chez un tiers US sans base légale.", "size": 13, "sa": 7},
    {"t": "Vendor lock-in : la douleur n°1 que personne ne certifie résoudre.", "size": 13, "sa": 7},
    {"t": "Coûts cloud / IA opaques : besoin d'un chiffrage honnête, sourcé, projetable.", "size": 13, "sa": 7},
    {"t": "Maturité technique : LLM ouverts, embeddings self-host, RAG → une base souveraine est enfin faisable hors GAFAM.", "size": 13, "sa": 4},
])
pagefoot(s, PAGE)
notes(s, """
Insister sur la distinction job ≠ solution (Clayton Christensen / Bob Moesta). Le dirigeant n'« embauche » pas Strate pour « designer son infra » — c'est le moyen. Il embauche Strate pour un RÉSULTAT souverain et une tranquillité d'esprit. (PRD §2.)

Détailler les jobs :
- FONCTIONNEL : prendre une décision d'infra engageante et difficilement réversible sans être expert, et pouvoir la maintenir dans le temps (prix/CGU/API qui bougent).
- ÉMOTIONNEL : ne pas avoir peur de se faire capturer / de payer trop / de perdre ses données ; sentiment de « contrôle souverain » (cf. DESIGN.md).
- SOCIAL : être perçu comme un dirigeant rigoureux qui maîtrise sa donnée et sa conformité (argument fort chez les régulés et le public).

Moments de douleur (triggers d'achat) : départ d'un sachant ; audit RGPD/CNDP ; hausse surprise de facture cloud ; changement de CGU d'un fournisseur US ; projet sur-mesure qui dérape.

Pourquoi maintenant : convergence réglementaire (RGPD chap. V sur les transferts, CNDP au Maroc, règle max(RGPD, CNDP) appliquée dans le produit) + maturité technique (embeddings open Apache 2.0 self-hostables type Qwen3-VL, RAG, LLM ouverts) qui rend une base souveraine faisable hors hyperscalers.

Ne PAS sur-vendre la taille du marché « souveraineté IA » : on ne dispose pas d'une stat de TAM sourcée — l'avocat du diable y reviendra (slide risques). Ici on reste sur la nature et la réalité de la douleur, pas sur le volume.
""")

# --- 4. SEGMENTS / PERSONAS (PMF 1/3) --------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Product-Market Fit · 1/3", "À qui on vend d'abord — et à qui pas encore", accent=BLUE, eb_color=BLUE)
# tableau personas en cartes empilées
rows = [
    ("P1 — Bâtisseur souverain", "Coach / freelance / consultant FR Maroc-UE, technophile, veut sa base revendable.", "CŒUR", TEAL),
    ("P2 — PME / startup tech", "Équipe 5-10, profil MEDIUM, cycle court, à l'aise avec ADR/specs/code.", "SECONDAIRE", BLUE),
    ("P3 — Agence / programme", "80+ accompagnés, multi-voix, volume + multi-tenancy.", "LOT ULTÉRIEUR", GOLD_TXT),
    ("P4 — Cabinet régulé", "Avocat / CGP / santé, secret pro, preset HARD, premium mais cycle long.", "LOT ULTÉRIEUR", GOLD_TXT),
    ("P5 — ETI / grand groupe", "Appel d'offres, DPA, SecNumCloud — incompatible self-service.", "HORS PÉRIMÈTRE", ERROR),
]
top = Inches(2.0); rh = Inches(0.84); gap = Inches(0.12); l = Inches(0.85); w = Inches(11.6)
for i, (name, desc, tag, col) in enumerate(rows):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    c = card(s, l, t, w, rh, fill=(WHITE if i < 2 else SURFACE_LOW))
    # bande couleur gauche
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.08), rh)
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    r1 = p.add_run(); set_run(r1, name + "   ", 15, INK, FONT_TITLE, bold=True)
    r2 = p.add_run(); set_run(r2, desc, 12.5, INK_VARIANT, FONT_BODY)
    # tag à droite
    tb2, tf2 = textbox(s, Emu(int(l) + int(w) - int(Inches(2.5))), t, Inches(2.35), rh, anchor=MSO_ANCHOR.MIDDLE)
    pp = tf2.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run(); set_run(rr, tag, 11.5, col, FONT_BODY, bold=True)
pagefoot(s, PAGE)
notes(s, """
La cible couvre 5 ordres de grandeur (Solo ~5 €/mois → Grand groupe ~1 M€/an) : ce ne sont NI le même acheteur, NI le même cycle de vente, NI la même conformité. La discipline ICP est donc cruciale. (PRD §3.)

Le wedge initial = P1 (Bâtisseur souverain), parce qu'Amine y a un accès direct (communauté), une demande quasi acquise et une preuve sociale / bouche-à-oreille. C'est le segment où l'on peut valider l'offre et le pricing sans payer un CAC élevé.

P2 (PME tech) = extension naturelle, cycle court, technophile.
P3 (agence/programme) et P4 (cabinet régulé) = lots ultérieurs : plus de valeur (multi-tenancy, premium conformité) mais cycle plus lourd / plus long.
P5 (ETI/grand groupe) = EXPLICITEMENT hors périmètre self-service (appel d'offres, DPA, SecNumCloud). On ne se disperse pas.

Risque connu (A2 du PRD) : « cible trop large ». Mitigation = focalisation P1/P2 d'abord. À répéter : la discipline de refus (dire non à P5 pour l'instant) est une décision, pas un aveu de faiblesse.

Note multi-tenant : les rails sont déjà posés (RLS Supabase, pivot circle) — donc passer à P3 plus tard est un choix de go-to-market, pas un chantier technique de fond.
""")

# --- 5. PMF 2/3 : Value Proposition Canvas ---------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Product-Market Fit · 2/3", "Value proposition : douleurs ↔ remèdes", accent=BLUE, eb_color=BLUE)
# deux colonnes : Pains (client) / Pain relievers (Strate)
c1 = card(s, Inches(0.85), Inches(2.0), Inches(5.7), Inches(4.6), fill=SURFACE_LOW)
card_text(c1, [
    {"t": "Douleurs du client", "size": 16, "color": ERROR, "font": FONT_TITLE, "bold": True, "sa": 10, "sb": 2},
    {"t": "Peur de se faire capturer (lock-in irréversible).", "size": 13, "sa": 8},
    {"t": "Données critiques chez un tiers US, hors conformité.", "size": 13, "sa": 8},
    {"t": "Coût d'infra opaque, factures qui dérivent.", "size": 13, "sa": 8},
    {"t": "Pas l'expertise pour arbitrer 7 couches techniques.", "size": 13, "sa": 8},
    {"t": "Conseil suspecté d'être commissionné par les vendors.", "size": 13, "sa": 4},
])
c2 = card(s, Inches(6.75), Inches(2.0), Inches(5.7), Inches(4.6), fill=WHITE)
card_text(c2, [
    {"t": "Remèdes Strate", "size": 16, "color": TEAL, "font": FONT_TITLE, "bold": True, "sa": 10, "sb": 2},
    {"t": "Exit Escrow : bundle reproductible, sortie en 1 clic.", "size": 13, "sa": 8},
    {"t": "Souveraineté progressive : UE / Maroc / on-prem, résidence à venir.", "size": 13, "sa": 8},
    {"t": "Coûts sourcés ±30 %, URL + date + confiance, feed live.", "size": 13, "sa": 8},
    {"t": "Reco 7 couches + ensemble multi-config + verdict 90 s.", "size": 13, "sa": 8},
    {"t": "Fiduciary Mode : zéro commission cachée, charte contractuelle.", "size": 13, "sa": 4},
])
pagefoot(s, PAGE)
notes(s, """
Value Proposition Canvas (Osterwalder) en version condensée : à gauche les douleurs (pains) du client, à droite les remèdes (pain relievers) RÉELLEMENT livrés par le produit.

Correspondance 1-pour-1, vérifiable dans le code :
- Lock-in → Exit Escrow (F7) : bundle IaC (Terraform/Docker Compose) + dumps DB + vault markdown source-de-vérité + script de re-embedding + runbook ; test d'intégration vérifie que le manifeste == reco. C'est la PREUVE DURE.
- Données hors conformité → souveraineté ; règle max(RGPD, CNDP) ; résidence/DR = épic à venir (S-043→048) : à présenter comme roadmap, pas comme livré.
- Coûts opaques → price feed live Firecrawl réconcilié à une baseline sourcée (garde-fou DÉFCON 1 : un prix aberrant est rejeté au profit de la baseline), chaque chiffre porte URL + date + pastille 🟢🟡🟠 + disclaimer « une IA peut se tromper ».
- Manque d'expertise → reco 7 couches + ensemble multi-config (montre l'incertitude) + verdict en 90 s ; intake en langage libre (LLM extrait, le serveur valide/borne) ; assistant Q&A contextuel.
- Conseil suspect → Fiduciary Mode (charte testée : le moteur ne connaît ni commission, ni affiliation, ni rétrocommission).

Le « gain creator » (côté positif) : repartir avec sa stack à tout moment = tranquillité. Mais ne pas chiffrer un gain de productivité sans mesure terrain (DÉFCON 1).

Hypothèses de PMF à valider (à dire à voix haute) : (1) P1 accepte de payer pour la « cuisine » et pas seulement de consommer la « recette » gratuite ; (2) l'Exit Escrow est un argument d'ACHAT et pas seulement de réassurance ; (3) la souveraineté est un critère de décision payant, pas un simple confort.
""")

# --- 6. PMF 3/3 : différenciation vs alternatives --------------------------
PAGE += 1
s = add_slide()
std_header(s, "Product-Market Fit · 3/3", "Face aux alternatives", accent=BLUE, eb_color=BLUE)
alts = [
    ("Statu quo (rien / ChatGPT)", "Mémoire faible, données chez un tiers US, aucune souveraineté.", "Strate : base souveraine + sortie certifiée."),
    ("Intégrateur / projet sur-mesure", "30-200 k€, 6 mois, non réutilisable, dépendance au prestataire.", "Strate : conseil chiffré ±30 %, reproductible, 90 s."),
    ("Hyperscalers (AWS/GCP/Azure)", "Puissants mais lock-in profond, complexité, opacité de coût.", "Strate : neutre, anti-lock-in, multi-fournisseurs."),
    ("RAG-as-a-Service (SaaS mémoire US)", "Locataire à vie, données hors UE, pas d'exit.", "Strate : exit garanti + fiduciaire + réseau de coûts."),
]
top = Inches(2.0); rh = Inches(1.06); gap = Inches(0.12); l = Inches(0.85); w = Inches(11.6)
for i, (name, weak, edge) in enumerate(alts):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    c = card(s, l, t, w, rh)
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    para(tf, name, 14, INK, FONT_TITLE, bold=True, first=True, space_after=3)
    p = tf.add_paragraph(); p.space_after = Pt(0)
    r1 = p.add_run(); set_run(r1, "Limite : ", 12, ERROR, FONT_BODY, bold=True)
    r2 = p.add_run(); set_run(r2, weak + "    ", 12, INK_VARIANT, FONT_BODY)
    r3 = p.add_run(); set_run(r3, "→ ", 12, TEAL, FONT_BODY, bold=True)
    r4 = p.add_run(); set_run(r4, edge, 12, TEAL, FONT_BODY, bold=True)
pagefoot(s, PAGE)
notes(s, """
Positionnement concurrentiel. Présenter chaque alternative honnêtement (l'avocat du diable challengera la défensibilité — voir slide risques).

- Statu quo / grand public : le concurrent réel le plus fréquent est « ne rien faire » ou ChatGPT/Claude. Faiblesse : mémoire faible, données US, zéro souveraineté.
- Intégrateur sur-mesure : 30-200 k€, 6 mois, non réutilisable (chiffres = repère de genèse P4, à traiter comme ordre de grandeur, pas comme stat sourcée). Strate offre un conseil chiffré, reproductible, en 90 s pour le diagnostic.
- Hyperscalers : ce sont AUSSI des fournisseurs potentiels DANS la stack recommandée — Strate n'est pas anti-cloud, il est anti-CAPTURE. Argument : neutralité + anti-lock-in.
- RAG-as-a-Service US : la cible directe à déloger ; faiblesse = locataire à vie, pas d'exit, données hors UE.

Le fil rouge de la différenciation : ce que les autres ne peuvent PAS faire sans saborder leur modèle, c'est garantir la sortie (Exit Escrow) et s'engager fiduciairement (zéro commission). Un acteur commissionné par les vendors ne peut pas copier le Fiduciary Mode — c'est la leçon Flipper (slide moats).

Caveat avocat du diable, à anticiper : un hyperscaler ou un SaaS pourrait COPIER l'Exit Escrow s'il y voyait un intérêt marketing. La vraie barrière n'est donc pas la feature seule, c'est la COMBINAISON (exit + fiduciaire + réseau de coûts) + l'incompatibilité de la copie avec leur modèle de revenu (lock-in). On y revient en section risques.
""")

# --- 7. LES 3 MOATS ---------------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Avantage défendable", "Le trio de moats (playbook Cloudflare)", accent=GOLD, eb_color=GOLD_TXT)
moats = [
    ("①  Exit Escrow", "Sortie certifiée", "Bundle reproductible (IaC + dumps + vault + runbook) redéployable ailleurs en 1 clic. Release auto même si Strate ferme.", "Analogie : escrow logiciel / IP Bankruptcy Act", TEAL, "PREUVE DURE · livré (F7)"),
    ("②  Fiduciary Mode", "Mandataire transparent", "Payé par l'utilisateur, jamais commissionné en douce par les vendors. La négociation tarifaire se fait au nom du client.", "Analogie : courtier énergie/assurance · leçon Flipper", GOLD_TXT, "Charte contractuelle · livré (F8)"),
    ("③  Intelligence Network", "Effet réseau de données", "Chaque déploiement renvoie (opt-in RGPD) le coût réel anonymisé → calibration collective irréplicable + alertes vendor.", "Analogie : ISAC cyber (collective defense)", BLUE, "Rails posés (F9) · actif au Lot 2"),
]
cw, ch, gap, top, l0 = Inches(3.7), Inches(4.45), Inches(0.25), Inches(2.0), Inches(0.85)
for i, (h, sub, body, ana, col, status) in enumerate(moats):
    l = Emu(int(l0) + i * (int(cw) + int(gap)))
    c = card(s, l, top, cw, ch, fill=WHITE)
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, top, cw, Inches(0.55))
    head.fill.solid(); head.fill.fore_color.rgb = col; head.line.fill.background(); head.shadow.inherit = False
    htf = head.text_frame; _no_autofit(htf); htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    htf.margin_left = Inches(0.16)
    hp = htf.paragraphs[0]; hr = hp.add_run(); set_run(hr, h, 16, WHITE, FONT_TITLE, bold=True)
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.72)
    para(tf, sub, 13.5, col, FONT_TITLE, bold=True, first=True, space_after=8)
    para(tf, body, 12.5, INK_VARIANT, space_after=10)
    para(tf, ana, 11, INK, FONT_BODY, italic=True, space_after=10)
    para(tf, status, 11, col, FONT_MONO, bold=True, space_after=0)
pagefoot(s, PAGE)
notes(s, """
Le cœur défendable. La couche de reco est commoditisée (un LLM la rapproche gratuitement). Le moat est le TRIO, qui transpose le playbook Cloudflare : se mettre devant l'origine sans capturer (Exit Escrow) ; créer l'habitude par la transparence (Fiduciary) ; voir tout le « trafic » de coûts pour mieux servir tous (Network). (MOAT-HUNT.md, scores 14/13/12 sur 15.)

① Exit Escrow (effort S, livré) : faisable PARCE QUE le vault markdown est déjà la source de vérité et que les projections (DB, vecteurs) sont des recalculs. Le bundle est testé : le manifeste doit correspondre exactement à la reco (couches/coût/preset). C'est la seule preuve DURE qu'on peut montrer aujourd'hui.

② Fiduciary Mode (effort S, livré) : LEÇON NÉGATIVE DÉCISIVE. Flipper (auto-switch énergie) a fermé ; Look After My Bills ne bascule que vers les vendors qui le commissionnent. Un courtier commissionné par le vendor TRAHIT le mandat et meurt. Donc le Fiduciary Mode n'est pas une option marketing, c'est une condition de survie ET une barrière : un concurrent commissionné ne peut pas copier sans saborder son modèle de revenu. (Sources : moneytothemasses, brabners — dans le PRD.)

③ Intelligence Network (effort L, rails posés) : c'est le VRAI moat Cloudflare (effet réseau de données). Chaque déploiement monitoré renvoie, sur opt-in RGPD horodaté, le coût réel par profil/vendor/région → dataset de calibration collectif qui tue progressivement le ±30 % + alertes type ISAC (un vendor change prix/CGU/API → tous alertés). MAIS : il ne devient ACTIF qu'avec des déploiements monitorés (Lot 2). Aujourd'hui = rails (consentement + schéma de données). À être honnête là-dessus (avocat du diable : effet réseau réel ? slide risques).

Sources des analogies toutes dans MOAT-HUNT.md (escrow logiciel, courtier énergie, ISAC) — ne jamais présenter ces analogies comme des garanties de succès, ce sont des modèles mentaux validés sur 10+ industries.
""")

# --- 8. GO-TO-MARKET --------------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Go-To-Market", "Du diagnostic 90 s au devis — motions & boucles", accent=TEAL)
# Funnel horizontal
steps = [
    ("Accroche", "Diagnostic gratuit\nintake libre / 4 questions", TEAL),
    ("Verdict 90 s", "Reco 7 couches + coût\nsourcé + risque/gain", TEAL_DIM),
    ("Livrable", "Export MD/PDF sourcé\n+ Exit Escrow", BLUE),
    ("Engagement", "Cuisine payante :\ndéploiement + supervision", GOLD_TXT),
]
fw, fh, gap, top, l0 = Inches(2.72), Inches(1.5), Inches(0.22), Inches(2.05), Inches(0.85)
for i, (h, body, col) in enumerate(steps):
    l = Emu(int(l0) + i * (int(fw) + int(gap)))
    c = card(s, l, top, fw, fh, fill=col, border=None)
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    para(tf, h, 15, WHITE, FONT_TITLE, bold=True, first=True, space_after=4, align=PP_ALIGN.CENTER)
    for line in body.split("\n"):
        para(tf, line, 11.5, WHITE, FONT_BODY, align=PP_ALIGN.CENTER, space_after=1)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Emu(int(l) + int(fw) - int(Inches(0.02))), Emu(int(top) + int(Inches(0.52))), Inches(0.26), Inches(0.46))
        ar.fill.solid(); ar.fill.fore_color.rgb = CARD_BORDER; ar.line.fill.background(); ar.shadow.inherit = False
# bandeau motions
tb, tf = textbox(s, Inches(0.85), Inches(3.95), Inches(11.6), Inches(2.9))
para(tf, "Motion & canaux", 16, INK, FONT_TITLE, bold=True, first=True, space_after=8)
bullet(tf, "Hybride PLG → sales-léger : le free tier (recette ouverte) crée l'habitude, l'humain ferme la « cuisine ».", bold_lead="Motion. ", space_after=8)
bullet(tf, "Communauté d'Amine (P1) en premier — accès direct, CAC quasi nul, bouche-à-oreille ; puis contenu + référencement.", bold_lead="Acquisition. ", space_after=8)
bullet(tf, "Exit Escrow = boucle de confiance (réduit la peur d'achat) ; Network = boucle de données (plus de clients → meilleure calibration → meilleur produit).", bold_lead="Boucles de croissance. ", space_after=8)
bullet(tf, "Maroc-UE francophone d'abord (P1/P2), puis régulés (P4) et agences (P3) ; ETI hors périmètre.", bold_lead="Géographie. ", space_after=0)
pagefoot(s, PAGE)
notes(s, """
Le chemin d'acquisition est déjà construit dans le produit : l'accueil propose un « chemin 90 s » (4 questions OU intake en langage libre / dictée) → /resultats en mode verdict (besoin / risque / gain mesuré / prix ferme [PLACEHOLDER] / coût infra ±30 % / mise en route / prochaine étape) → livrable exportable (MD/PDF) + bundle Exit Escrow → upsell « cuisine payante ».

Motion :
- PLG (Product-Led Growth) en haut de funnel : le diagnostic gratuit est le lead magnet (modèle Cloudflare free tier). Faible friction, valeur immédiate, capture d'e-mail exit-intent (RLS, opt-in RGPD/CNDP).
- Sales-léger en bas de funnel : la conversion vers la cuisine payante (déploiement assisté, supervision) demande un échange humain, surtout sur les segments à enjeu (régulés). On n'a pas besoin d'une force de vente lourde au début (P1 = communauté).

Boucles de croissance :
1. Boucle de confiance (moat ① + ②) : Exit Escrow + charte fiduciaire réduisent la peur d'achat → plus de conversions → plus de preuves sociales.
2. Boucle de données (moat ③) : chaque déploiement monitoré enrichit la calibration des coûts → meilleur produit → meilleure conversion. Cette boucle ne tourne qu'au Lot 2 (déploiements monitorés) — honnêteté.

Séquencement géo : francophone Maroc-UE d'abord (là où Amine a l'accès et où max(RGPD, CNDP) est un argument), puis élargissement.

Métriques de funnel cibles (PRD §10, à traiter comme cibles, pas résultats) : activation > 60 % (complètent le wizard → livrable) ; conversion P1 livrable → payant > 10 % ; rétention (comptes ayant généré un Exit bundle) > 80 %. Ce sont des HYPOTHÈSES de cible.
""")

# --- 9. PRICING 1/2 : architecture -----------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Stratégie de prix · 1/2", "Architecture : on tarife le service, pas l'infra", accent=GOLD, eb_color=GOLD_TXT)
# principe central
c0 = card(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(1.0), fill=SURFACE_LOW)
card_text(c0, [
    {"t": "Le service Strate = prix ferme.  L'infra = coût aux fournisseurs (±30 %, zéro marge).  « Recette ouverte, cuisine payante. »",
     "size": 15.5, "color": TEAL, "font": FONT_TITLE, "bold": True, "sa": 0, "align": PP_ALIGN.CENTER},
], anchor=MSO_ANCHOR.MIDDLE)
# 3 paliers candidats
tiers = [
    ("Diagnostic", "Gratuit", "Wizard + reco + livrable de base sourcé. Lead magnet, crée l'habitude.", CARD_BORDER, INK),
    ("Pro", "[PLACEHOLDER] / mois", "Déploiement assisté + monitoring + recalibration + négociation vendor. Exit bundle premium.", TEAL, WHITE),
    ("Souverain+", "Sur devis", "On-prem / air-gapped, gouvernance avancée, support dédié + SLA, multi-région (à venir).", GOLD_TXT, WHITE),
]
cw, ch, gap, top, l0 = Inches(3.7), Inches(3.0), Inches(0.25), Inches(3.25), Inches(0.85)
for i, (name, price, body, col, txtcol) in enumerate(tiers):
    l = Emu(int(l0) + i * (int(cw) + int(gap)))
    featured = (i == 1)
    c = card(s, l, top, cw, ch, fill=(col if i > 0 else WHITE), border=(TEAL if featured else CARD_BORDER), border_w=(2.2 if featured else 1.0))
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.22)
    name_col = WHITE if i > 0 else INK
    para(tf, name, 18, name_col, FONT_TITLE, bold=True, first=True, space_after=4, align=PP_ALIGN.CENTER)
    price_col = WHITE if i > 0 else TEAL
    para(tf, price, 20, price_col, FONT_MONO, bold=True, space_after=10, align=PP_ALIGN.CENTER)
    body_col = RGBColor(0xEE, 0xF0, 0xFF) if i > 0 else INK_VARIANT
    para(tf, body, 12.5, body_col, space_after=0, align=PP_ALIGN.CENTER)
pagefoot(s, PAGE)
notes(s, """
PRINCIPE D'ARCHITECTURE (PRD §9, charte fiduciaire) : on ne marge JAMAIS sur l'infra. Le coût cloud/IA est transparent, payé directement par le client à ses fournisseurs (±30 %, sourcé). Strate facture uniquement le SERVICE : déploiement assisté, monitoring, recalibration, négociation tarifaire vendor au nom du client (rémunération affichée). « Recette ouverte, cuisine payante. »

Cette séparation est ce qui rend le Fiduciary Mode crédible : si on margeait sur l'infra, on aurait intérêt à pousser des composants chers — exactement le conflit qu'on dénonce.

Les 3 paliers (grille candidate, cf. docs/pricing/wtp-research.md §3) :
- Diagnostic = GRATUIT. Lead magnet. Le QUOI est offert (anti-lock-in + habitude).
- Pro = forfait mensuel = [PLACEHOLDER]. C'est LE prix à caler par le sondage. Contenu : déploiement assisté + supervision + recalibration + négo vendor + Exit bundle premium.
- Souverain+ = sur devis. On-prem/air-gapped, gouvernance, SLA, multi-région (résidence/DR = épic à venir).

IMPORTANT (DÉFCON 1) : le prix « Pro » n'est PAS encore connu. La homepage affiche un badge « Prix, sondage en cours » et ne donne aucun montant « au doigt mouillé ». Ne jamais asséner un chiffre ici.

Engagement : mensuel ou annuel (−20 % en hypothèse, attribut testé dans le conjoint). L'annuel finance le runway et baisse le churn.
""")

# --- 10. PRICING 2/2 : WTP + fourchettes -----------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Stratégie de prix · 2/2", "Caler le prix : la méthode, pas le doigt mouillé", accent=GOLD, eb_color=GOLD_TXT)
c1 = card(s, Inches(0.85), Inches(2.0), Inches(5.7), Inches(4.55), fill=WHITE)
card_text(c1, [
    {"t": "Méthode (en cours)", "size": 16, "color": TEAL, "font": FONT_TITLE, "bold": True, "sa": 9, "sb": 2},
    {"t": "Van Westendorp (10-15 prospects) → fourchette acceptable [PMC … PME], OPP, IPP.", "size": 12.5, "sa": 8},
    {"t": "Conjoint CBC (30-50 répondants) → quels attributs portent le prix + prix optimal par palier.", "size": 12.5, "sa": 8},
    {"t": "Question d'intention (Q5) → corrige l'optimisme déclaratif.", "size": 12.5, "sa": 8},
    {"t": "Seuil de figeage : ≥ 15 réponses VW + ≥ 30 conjoint avant de fixer.", "size": 12.5, "sa": 2},
])
c2 = card(s, Inches(6.75), Inches(2.0), Inches(5.7), Inches(4.55), fill=SURFACE_LOW)
card_text(c2, [
    {"t": "Modèles de revenu candidats", "size": 16, "color": GOLD_TXT, "font": FONT_TITLE, "bold": True, "sa": 9, "sb": 2},
    {"t": "Abonnement (forfait service /mois) — base, prévisible.", "size": 12.5, "sa": 8},
    {"t": "Setup + run (mise en route one-time + récurrent).", "size": 12.5, "sa": 8},
    {"t": "À l'usage / value-based (négo vendor : part de l'économie générée).", "size": 12.5, "sa": 8},
    {"t": "Attributs de valeur (conjoint) : on-prem, recalibration/négo, support dédié.", "size": 12.5, "sa": 2},
])
# bandeau bas
tb, tf = textbox(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.4))
para(tf, "Prix de vente Strate = [PLACEHOLDER] tant que le sondage n'a pas atteint ses seuils. Aucune fourchette ferme présentée ici — c'est une décision à trancher après données.",
     11, ERROR, FONT_MONO, bold=False, first=True)
pagefoot(s, PAGE)
notes(s, """
Comment on calera le prix (docs/pricing/wtp-research.md) — la méthode est le livrable, pas un chiffre :

1. Mini-test Van Westendorp (4 questions de prix : trop cher / trop bon marché / cher mais acceptable / bonne affaire) sur 10-15 prospects P1. Donne : OPP (Optimal Price Point = intersection trop cher × trop bon marché), IPP (Indifference Price Point ≈ prix médian psychologique), et la fourchette acceptable [PMC … PME]. Une question d'intention (Q5) corrige l'optimisme du déclaratif.

2. Conjoint Choice-Based (CBC) sur 30-50 répondants : attributs = déploiement (UE / hybride / on-prem air-gapped) ; périmètre (déploiement seul / + monitoring / + recalibration & négo) ; capacités (base / + gouvernance / + création) ; support (standard / prioritaire / dédié + SLA) ; engagement (mensuel / annuel −20 %) ; prix (4 niveaux). Sortie : utilités partielles, importance relative des attributs, WTP par niveau, simulateur de parts → prix optimal PAR palier.

3. Seuil de figeage : NE PAS fixer le prix avant ≥ 15 réponses VW et ≥ 30 conjoint. (Discipline DÉFCON 1.)

Modèles de revenu candidats (à arbitrer une fois la WTP connue) :
- Abonnement forfaitaire (le plus simple, prévisible) — probablement le socle.
- Setup + run (capture la valeur de la mise en route, qui est réelle : ingestion du backlog, configuration).
- Value-based sur la négo vendor : facturer une part de l'économie obtenue au nom du client (aligné fiduciaire SI affiché).

Attributs qui portent probablement le prix (hypothèse à confirmer par le conjoint) : on-prem/air-gapped, recalibration + négo vendor, support dédié + SLA.

Le message clé pour Amine : on présente une MÉTHODE rigoureuse et un dispositif déjà en place (page /sondage + API de collecte), pas un prix inventé. Le [PLACEHOLDER] est assumé et c'est une force de crédibilité, pas une faiblesse.
""")

# --- 11. COMMUNICATION ------------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Stratégie de communication", "Un récit de souveraineté, prouvé par l'exit", accent=BLUE, eb_color=BLUE)
tb, tf = textbox(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(0.85))
para(tf, "« Votre mémoire d'organisation, souveraine et réversible. Repartez avec toute votre stack quand vous voulez. »",
     18, BLUE, FONT_TITLE, bold=True, italic=True, first=True)
# 3 colonnes : récit / preuves / messages par segment
cols = [
    ("Récit de marque", ["Souveraineté (vos données, votre loi).", "Fiduciaire (zéro commission cachée).", "Anti-lock-in (sortie certifiée).", "Transparence radicale des coûts."], TEAL),
    ("Régime de preuve", ["DURE : Exit Escrow (testable) + charte fiduciaire.", "MOLLE : gain de temps / qualité = « à mesurer chez vous ».", "Jamais de stat inventée (DÉFCON 1)."], GOLD_TXT),
    ("Messages par segment", ["P1 : votre base, revendable, à vous.", "P2 : conformité + vitesse, sans intégrateur.", "P4 : secret pro, audit, preset HARD."], BLUE),
]
cw, ch, gap, top, l0 = Inches(3.7), Inches(3.35), Inches(0.25), Inches(3.0), Inches(0.85)
for i, (h, items, col) in enumerate(cols):
    l = Emu(int(l0) + i * (int(cw) + int(gap)))
    c = card(s, l, top, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, top, cw, Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf)
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.24)
    para(tf, h, 15, INK, FONT_TITLE, bold=True, first=True, space_after=9)
    for it in items:
        bullet(tf, it, size=12, marker_color=col, space_after=8)
pagefoot(s, PAGE)
notes(s, """
Le récit de marque s'appuie sur l'identité produit (DESIGN.md : « sovereign control », crédibilité d'ingénierie, transparence, stabilité premium). Trois piliers narratifs : souveraineté, fiduciaire, anti-lock-in.

RÉGIME DE PREUVE (point capital, DÉFCON 1) : la homepage RÉELLE distingue déjà deux registres —
- PREUVE DURE : Exit Escrow (on peut le tester : le bundle redéploie une stack identique) + charte fiduciaire (engagement contractuel vérifiable). Ce sont les seules promesses « démontrables » qu'on affiche comme telles.
- PREUVE MOLLE : tout ce qui touche au gain de temps, à la qualité de réponse, à la productivité = « à mesurer chez vous ». On NE chiffre PAS un bénéfice sans donnée client. Zéro stat inventée.

Cette honnêteté EST un message marketing : dans un marché saturé de promesses IA gonflées, dire « voici ce qu'on peut prouver, voici ce qu'il faudra mesurer » construit la confiance — cohérent avec le positionnement fiduciaire.

Messages par segment :
- P1 (bâtisseur) : « votre base, à vous, revendable » — appropriation + valeur d'actif.
- P2 (PME tech) : « conformité + vitesse, sans payer un intégrateur 6 mois ».
- P4 (régulé) : « secret professionnel, audit, bi-temporalité, preset HARD » — conformité comme argument premium.

Canaux : contenu (la rigueur du chiffrage sourcé = différenciant éditorial), communauté d'Amine, démonstrations live de l'Exit Escrow (le « moment wow » : générer un bundle et le redéployer ailleurs). Positionnement vs concurrents : « eux vous gardent locataire, nous vous rendons propriétaire ».

Garde-fou : la voix de marque doit rester non-orientée (principe « avis critique non orienté » d'Amine) — Strate présente alternatives/fourchettes/sources, le client tranche. Ne pas glisser vers un discours directif.
""")

# --- 12. P&L INTRO / drivers -----------------------------------------------
PAGE += 1
s = add_slide(INVERSE)
accent_bar(s, GOLD)
eyebrow(s, "P&L prévisionnel · 1/3", RGBColor(0xFF, 0xBA, 0x46))
title(s, "Modèle financier — tout est HYPOTHÈSE", color=WHITE)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.78), Inches(11.6), Pt(1.4))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x4a,0x53,0x66); ln.line.fill.background(); ln.shadow.inherit=False
tb, tf = textbox(s, Inches(0.85), Inches(2.1), Inches(11.6), Inches(0.9))
para(tf, "Aucun chiffre des 3 prochaines slides n'est une prévision. Ce sont des DRIVERS et des scénarios pour raisonner — à valider par le sondage et les premières ventes.",
     15, RGBColor(0xFF, 0xBA, 0x46), FONT_BODY, italic=True, first=True)
# drivers en 2 colonnes
tb, tf = textbox(s, Inches(0.85), Inches(3.15), Inches(5.7), Inches(3.6))
para(tf, "Drivers de revenu", 15, RGBColor(0x6B,0xD8,0xCB), FONT_TITLE, bold=True, first=True, space_after=9)
for d in ["Nombre de clients payants (Pro + Souverain+).", "ARPA = prix service /mois [PLACEHOLDER].", "Mix paliers (Pro vs Souverain+ devis).", "Taux de conversion diagnostic → payant.", "Churn mensuel / rétention."]:
    bullet(tf, d, size=12.5, color=INVERSE_TXT, marker_color=RGBColor(0x6B,0xD8,0xCB), space_after=7)
tb, tf = textbox(s, Inches(6.75), Inches(3.15), Inches(5.7), Inches(3.6))
para(tf, "Drivers de coût", 15, RGBColor(0xFF,0xBA,0x46), FONT_TITLE, bold=True, first=True, space_after=9)
for d in ["Marge brute : l'infra est REVERSÉE (pas de marge) → revenu = service.", "Coûts LLM/tokens à l'usage (intake, narration, assistant, veille).", "Firecrawl (price feed + veille catalogue).", "CAC (faible sur P1, croissant ensuite).", "OPEX : dev, support, hébergement plateforme."]:
    bullet(tf, d, size=12.5, color=INVERSE_TXT, marker_color=RGBColor(0xFF,0xBA,0x46), space_after=7)
pagefoot(s, PAGE)
notes(s, """
Slide de cadrage du P&L, sur fond sombre pour marquer la rupture : ON ENTRE DANS LES HYPOTHÈSES. À dire explicitement : aucun chiffre des 3 slides suivantes n'est une prévision engageante. Ce sont des structures de raisonnement.

Drivers de REVENU :
- Nombre de clients payants : la variable maîtresse.
- ARPA (Average Revenue Per Account) = prix du service /mois = [PLACEHOLDER]. Tant qu'il n'est pas calé, tout chiffre d'affaires est paramétrique.
- Mix de paliers : combien de Pro (forfait) vs Souverain+ (devis, ARPA bien plus élevé mais cycle plus long).
- Conversion diagnostic → payant : cible HYPOTHÉTIQUE > 10 % sur P1.
- Churn / rétention : cible HYPOTHÉTIQUE rétention > 80 % (les comptes ayant généré un Exit bundle).

Drivers de COÛT — point structurant et CONTRE-INTUITIF :
- La marge brute n'inclut PAS l'infra, car l'infra est reversée au fournisseur SANS marge (fiduciaire). Donc le revenu Strate ≈ revenu de service quasi « pur ». La marge brute sur le SERVICE est potentiellement élevée (c'est du logiciel + de l'expertise), MAIS attention aux coûts variables d'IA.
- Coûts LLM/tokens : le produit appelle un LLM (via proxy LiteLLM/DeepSeek) pour l'intake libre, la narration, l'assistant Q&A et la veille catalogue ; + Firecrawl pour le price feed et la veille. Ces coûts sont VARIABLES et à l'usage — c'est une dette transparente notée dans le projet (le coût LLM à l'usage est signalé NON chiffré dans le moteur, story future). Ils grèvent la marge si le free tier est très utilisé → à surveiller (avocat du diable).
- CAC : faible au début (communauté P1, bouche-à-oreille), mais montera quand il faudra acquérir hors réseau.
- OPEX : développement (équipe réduite aujourd'hui), support, hébergement de la plateforme (Coolify/serveur).

Le vrai risque économique n'est pas « est-ce rentable par client » (le service marge bien) mais « le free tier coûte-t-il trop cher en IA avant conversion ? » et « le CAC hors-communauté est-il soutenable ? ». On y revient en risques.
""")

# --- 13. P&L 3 scénarios ----------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "P&L prévisionnel · 2/3", "Trois scénarios paramétriques (hypothèses)", accent=GOLD, eb_color=GOLD_TXT)
# tableau 3 scénarios x 3 ans
headers = ["", "An 1", "An 2", "An 3"]
data = [
    ("Prudent — clients payants", "8", "25", "60"),
    ("Base — clients payants", "15", "60", "180"),
    ("Ambitieux — clients payants", "30", "150", "500"),
]
# table via shapes
tx, ty, tw = Inches(0.85), Inches(2.15), Inches(11.6)
col_w = [Inches(5.0), Inches(2.2), Inches(2.2), Inches(2.2)]
rh = Inches(0.62)
# header row
x = int(tx)
for j, htxt in enumerate(headers):
    cell = card(s, Emu(x), ty, col_w[j], rh, fill=INVERSE, border=None, radius=False)
    tf = cell.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14)
    p = tf.paragraphs[0]; p.alignment = (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    r = p.add_run(); set_run(r, htxt, 13, WHITE, FONT_TITLE, bold=True)
    x += int(col_w[j])
# data rows
rowcolors = [SURFACE_LOW, WHITE, SURFACE_LOW]
rowaccent = [INK_VARIANT, TEAL, GOLD_TXT]
for ri, (name, a1, a2, a3) in enumerate(data):
    y = int(ty) + int(rh) * (ri + 1)
    x = int(tx)
    vals = [name, a1, a2, a3]
    for j, v in enumerate(vals):
        cell = card(s, Emu(x), Emu(y), col_w[j], rh, fill=rowcolors[ri], border=CARD_BORDER, border_w=0.75, radius=False)
        tf = cell.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        p = tf.paragraphs[0]; p.alignment = (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        r = p.add_run()
        set_run(r, v, (13 if j == 0 else 15), (rowaccent[ri] if j == 0 else INK),
                (FONT_BODY if j == 0 else FONT_MONO), bold=(j == 0 or j > 0))
        x += int(col_w[j])
# CA formule
tb, tf = textbox(s, Inches(0.85), Inches(4.75), Inches(11.6), Inches(2.0))
para(tf, "Chiffre d'affaires = clients payants × ARPA × 12, où ARPA = [PLACEHOLDER]/mois (sondage en cours).", 14, INK, FONT_MONO, bold=True, first=True, space_after=10)
bullet(tf, "Point mort : atteint quand CA service > OPEX fixes + coûts variables IA. Dépend entièrement de l'ARPA → calculable seulement après le sondage.", bold_lead="Seuil. ", space_after=8)
bullet(tf, "Runway : fonction de la base de coûts (équipe réduite, hébergement maîtrisé) — l'annuel (−20 %) sécurise la trésorerie.", bold_lead="Trésorerie. ", space_after=0)
pagefoot(s, PAGE)
notes(s, """
TROIS SCÉNARIOS — uniquement le driver maître (nombre de clients payants) est chiffré, et c'est une HYPOTHÈSE de modélisation, pas une prévision. On ne chiffre PAS le CA en euros parce que l'ARPA est un [PLACEHOLDER] : poser un CA fermerait une donnée non connue (interdit DÉFCON 1).

Logique des scénarios (à expliquer comme des « et si », pas des promesses) :
- PRUDENT (8 → 25 → 60) : adoption lente, conversion sous la cible, surtout P1, peu de Souverain+. C'est le scénario de survie.
- BASE (15 → 60 → 180) : la communauté P1 convertit + extension P2, quelques Souverain+ devis. Croissance organique saine.
- AMBITIEUX (30 → 150 → 500) : effet réseau qui s'enclenche (moat ③ actif), bouche-à-oreille fort, ouverture P3/P4. Suppose un Lot 2 mûr et un CAC maîtrisé.

Comment lire : multiplier par ARPA × 12 pour le CA. Exemple de raisonnement à faire AVEC un prix une fois connu — ne PAS donner de chiffre en euros maintenant.

Point mort : déterminé par CA service vs (OPEX fixes + coûts variables IA). Comme l'ARPA n'est pas fixé, le point mort est paramétrique. Une fois le sondage clos, on instancie 3 prix candidats et on lit le point mort de chacun.

Runway : la base de coûts est aujourd'hui légère (équipe réduite, hébergement self-host Coolify/serveur maîtrisé, pas de force de vente lourde). L'engagement annuel (−20 %) améliore la prévisibilité de trésorerie.

Le risque caché ici : les COÛTS VARIABLES IA du free tier. Si beaucoup d'utilisateurs font tourner intake + assistant + veille sans convertir, le coût LLM/Firecrawl par lead peut éroder la marge. Mitigation possible : caps/quotas sur le free tier, modèles moins chers via LiteLLM, cache. À traiter (slide annexe hypothèses + risques).

À répéter : ces nombres servent à structurer la discussion d'investissement et de runway, pas à promettre un résultat.
""")

# --- 14. P&L unit economics ------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "P&L prévisionnel · 3/3", "Unit economics : où se joue la viabilité", accent=GOLD, eb_color=GOLD_TXT)
ue = [
    ("Marge sur service", "Élevée en théorie (logiciel + expertise). L'infra reversée ne pèse PAS sur la marge Strate.", TEAL),
    ("Coût variable IA / lead", "Risque réel : LLM (intake, assistant, veille) + Firecrawl par utilisateur, même gratuit.", ERROR),
    ("CAC", "Quasi nul sur P1 (communauté) ; à modéliser dès qu'on acquiert hors réseau.", BLUE),
    ("LTV", "Fonction de l'ARPA [PLACEHOLDER] × durée de vie (rétention Exit-Escrow > 80 % visée).", GOLD_TXT),
]
cw, ch, gap, top, l0 = Inches(2.78), Inches(2.75), Inches(0.18), Inches(2.1), Inches(0.85)
for i, (h, body, col) in enumerate(ue):
    l = Emu(int(l0) + i * (int(cw) + int(gap)))
    c = card(s, l, top, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, top, cw, Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    card_text(c, [
        {"t": h, "size": 15, "color": INK, "font": FONT_TITLE, "bold": True, "sa": 8, "sb": 6},
        {"t": body, "size": 12.5, "color": INK_VARIANT, "sa": 4},
    ])
c0 = card(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(1.4), fill=SURFACE_LOW)
card_text(c0, [
    {"t": "Verrou de viabilité : LTV/CAC sain SI (1) le free tier IA reste sous contrôle (quotas / cache / modèle économe) et (2) la rétention tient grâce à l'Exit Escrow + la supervision continue.",
     "size": 14, "color": INK, "bold": False, "sa": 5},
    {"t": "Tout ratio chiffré ici serait fictif tant que ARPA, CAC et churn ne sont pas observés. → priorité : sondage + 5 premières ventes P1.",
     "size": 13, "color": GOLD_TXT, "bold": True, "italic": True, "sa": 0},
], anchor=MSO_ANCHOR.MIDDLE)
pagefoot(s, PAGE)
notes(s, """
Les unit economics sont LE sujet que l'avocat du diable doit forcer. On expose les 4 briques sans inventer de ratio (on n'a pas les données).

1. Marge sur service : potentiellement élevée. Comme l'infra est reversée sans marge (fiduciaire), elle ne grève PAS la marge de Strate — c'est un coût du CLIENT, pas de Strate. Donc le revenu Strate est presque entièrement du service à forte marge brute logicielle. C'est une bonne nouvelle structurelle.

2. Coût variable IA par lead : le PIÈGE. Même un utilisateur gratuit déclenche des appels LLM (intake libre, narration, assistant Q&A) et Firecrawl (price feed, veille catalogue). Ce coût existe AVANT toute conversion. Si le ratio leads gratuits / conversions est élevé, le coût d'acquisition « caché » grimpe. Mitigations : quotas sur le free tier, cache (déjà en place pour les prix et le catalogue, TTL court), bascule sur un modèle moins cher via LiteLLM (le provider se règle côté proxy, l'app référence un alias).

3. CAC : ~0 sur P1 (communauté, bouche-à-oreille). Mais on ne peut pas extrapoler ce CAC aux segments hors réseau — il faudra le mesurer (contenu, SEO, démos).

4. LTV : ARPA [PLACEHOLDER] × durée de vie. La rétention visée (>80 %) repose sur deux ancres : l'Exit Escrow (paradoxe : pouvoir partir facilement RETIENT, car ça enlève la peur) et la supervision continue (valeur récurrente). Hypothèse à valider.

Conclusion honnête : on ne PEUT PAS donner un LTV/CAC chiffré crédible aujourd'hui. Le faire serait du faux-rigoureux. La priorité est donc : (a) clore le sondage WTP, (b) réaliser 5 ventes P1 réelles pour observer conversion, ARPA effectif, churn et coût IA/lead. C'est le kill-criteria économique.
""")

# --- 15. PARTENARIATS -------------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Partenariats & alignement stratégique", "Nouer des alliances sans trahir le fiduciaire", accent=TEAL)
# tension box
c0 = card(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(0.95), fill=RGBColor(0xFF, 0xDA, 0xD6))
card_text(c0, [
    {"t": "Tension centrale : tout partenariat commissionné par un fournisseur DÉTRUIT le moat fiduciaire (leçon Flipper). Règle : aucune reco conditionnée à une rémunération vendor — jamais.",
     "size": 14.5, "color": RGBColor(0x93, 0x00, 0x0a), "bold": True, "sa": 0, "align": PP_ALIGN.CENTER},
], anchor=MSO_ANCHOR.MIDDLE)
parts = [
    ("Fournisseurs d'infra souveraine", "Scaleway, OVH, hébergeurs UE/Maroc. Accord : tarif négocié AU NOM du client, rémunération AFFICHÉE — jamais une rétrocommission cachée.", TEAL),
    ("Intégrateurs & cabinets conseil", "Sous-traitance du déploiement Lot 2 / apport d'affaires transparent. Strate reste le mandataire du client.", BLUE),
    ("Cabinets juridiques / DPO", "Couche conformité (RGPD/CNDP) en produit dérivé ; Strate reste « ingénierie, pas conseil juridique ».", GOLD_TXT),
    ("Open-source & fournisseurs de modèles", "Communautés (embeddings open, LLM) + proxy LiteLLM neutre : on change de modèle sans changer de produit.", TEAL),
]
cw, ch, gap, top, l0 = Inches(5.7), Inches(1.65), Inches(0.2), Inches(3.2), Inches(0.85)
for i, (h, body, col) in enumerate(parts):
    row, coln = divmod(i, 2)
    l = Emu(int(l0) + coln * (int(cw) + int(gap)))
    t = Emu(int(top) + row * (int(ch) + int(Inches(0.18))))
    c = card(s, l, t, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.08), ch)
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18)
    para(tf, h, 14, INK, FONT_TITLE, bold=True, first=True, space_after=4)
    para(tf, body, 12, INK_VARIANT, space_after=0)
pagefoot(s, PAGE)
notes(s, """
La question piège des partenariats : comment monétiser des alliances SANS détruire le moat fiduciaire (zéro commission cachée) ? C'est une vraie tension, pas un détail.

Règle d'or (charte fiduciaire, testée dans le code) : aucune recommandation n'est jamais conditionnée à une rémunération vendor. Le moteur de reco ne connaît NI commission, NI affiliation, NI rétrocommission. La leçon Flipper est explicite : un courtier commissionné par le vendor trahit le mandat et meurt.

Comment nouer des partenariats malgré tout :
1. Fournisseurs d'infra souveraine (Scaleway, OVH, hébergeurs UE/Maroc) : le Fiduciary Mode payant prévoit justement la NÉGOCIATION tarifaire au nom du client, avec rémunération AFFICHÉE. La nuance fiduciaire : on peut être rémunéré, mais de façon transparente et divulguée, jamais en douce, et sans que ça oriente la reco. Un rabais négocié reversé au client (ou une rémunération de service affichée) est compatible ; une rétrocommission cachée ne l'est pas.
2. Intégrateurs & cabinets conseil : pour exécuter le Lot 2 (déploiement) à l'échelle, Strate peut sous-traiter — en restant le mandataire du client, avec apport d'affaires transparent.
3. Cabinets juridiques / DPO : la conformité (RGPD/CNDP) est un produit dérivé naturel (Lot 3). Garde-fou : Strate reste « ingénierie, pas conseil juridique » (disclaimer présent dans la roadmap résidence/DR). Le partenariat juridique COMBLE ce vide sans que Strate s'expose au risque de practice du droit.
4. Open-source & fournisseurs de modèles : le proxy LiteLLM rend Strate neutre vis-à-vis des modèles (on change de provider côté proxy sans toucher le produit). Aligné avec la souveraineté et l'anti-dépendance. Les communautés open (embeddings Apache 2.0 type Qwen3-VL) renforcent le récit souverain.

Alignement avec la mission : tout partenariat doit servir l'un des 3 moats (souveraineté/anti-lock-in, transparence, données) — pas les diluer. Un partenariat qui pousserait un fournisseur pour une commission serait un anti-pattern à refuser, même lucratif.
""")

# --- 16. AVOCAT DU DIABLE — intro -------------------------------------------
PAGE += 1
s = add_slide(INVERSE)
accent_bar(s, ERROR)
eyebrow(s, "Avocat du diable", RGBColor(0xFF, 0xB4, 0xAB))
title(s, "Ce qui pourrait tuer Strate", color=WHITE)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.78), Inches(11.6), Pt(1.4))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x4a,0x53,0x66); ln.line.fill.background(); ln.shadow.inherit=False
tb, tf = textbox(s, Inches(0.85), Inches(2.15), Inches(11.6), Inches(4.5))
para(tf, "Sans complaisance — chaque risque appelle une expérience ou une preuve, pas une réassurance.", 15, RGBColor(0xFF, 0xBA, 0x46), FONT_BODY, italic=True, first=True, space_after=16)
risks_intro = [
    "Unit economics réelles : la marge survit-elle au coût IA du free tier et au CAC hors-communauté ?",
    "Dépendance fournisseurs (LLM / Firecrawl) : prêcher la souveraineté en dépendant d'API tierces ?",
    "Érosion des moats : un hyperscaler peut-il copier l'Exit Escrow ? l'effet réseau est-il réel ?",
    "Risque juridique : Strate frôle-t-il le conseil réglementé (juridique / financier) ?",
    "Marché : la « souveraineté IA » est-elle un marché assez grand et mûr, et payant ?",
    "Le piège du [PLACEHOLDER] : et si la willingness-to-pay est trop basse pour la cuisine ?",
]
for r in risks_intro:
    bullet(tf, r, size=14, color=INVERSE_TXT, marker="▸", marker_color=RGBColor(0xFF,0xB4,0xAB), space_after=11)
pagefoot(s, PAGE)
notes(s, """
Slide de transition vers la section critique. Ton : lucide, pas défaitiste. Le but est de transformer chaque angle aveugle en EXPÉRIENCE à mener (kill-criteria), conformément à la culture « zéro complaisance » d'Amine.

Les 6 axes annoncés ici sont développés sur les 2 slides suivantes :
1. Unit economics réelles (le coût IA caché + CAC).
2. Dépendance aux fournisseurs (paradoxe souveraineté vs dépendance API).
3. Érosion des moats (copiabilité de l'Exit Escrow, réalité de l'effet réseau).
4. Risque juridique (frontière conseil/ingénierie).
5. Taille et maturité du marché « souveraineté IA » (PAS de TAM sourcé → c'est une inconnue, pas un acquis).
6. Le piège du prix [PLACEHOLDER] (WTP potentiellement basse).

À garder en tête : l'absence de stat de marché sourcée est elle-même un risque (on ne peut pas affirmer que le marché est grand). On l'assume.
""")

# --- 17. AVOCAT DU DIABLE — risques détaillés (table 1) --------------------
PAGE += 1
s = add_slide()
std_header(s, "Avocat du diable · risques 1/2", "Risque → angle aveugle → preuve à mener", accent=ERROR, eb_color=ERROR)
risk_rows = [
    ("Unit economics", "Coût LLM/Firecrawl du free tier + CAC hors-communauté peuvent éroder la marge avant conversion.", "Mesurer coût IA/lead + conversion sur 50 diagnostics réels ; tester quotas/cache."),
    ("Dépendance fournisseurs", "Souveraineté prêchée mais intake/narration/assistant dépendent d'un LLM tiers ; price feed dépend de Firecrawl.", "Repli déterministe déjà codé ; tester un LLM auto-hébergé via LiteLLM ; dégradation gracieuse vérifiée."),
    ("Érosion des moats", "L'Exit Escrow est copiable ; l'effet réseau n'existe qu'avec un volume de déploiements monitorés.", "Mesurer le temps avant 1er effet réseau ; documenter pourquoi un commissionné ne peut pas copier le fiduciaire."),
]
top = Inches(2.0); rh = Inches(1.45); gap = Inches(0.16); l = Inches(0.85); w = Inches(11.6)
for i, (name, blind, exp) in enumerate(risk_rows):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    c = card(s, l, t, w, rh)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.09), rh)
    bar.fill.solid(); bar.fill.fore_color.rgb = ERROR; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    para(tf, name, 15, ERROR, FONT_TITLE, bold=True, first=True, space_after=4)
    p = tf.add_paragraph(); p.space_after = Pt(4)
    r1 = p.add_run(); set_run(r1, "Angle aveugle : ", 12, INK, FONT_BODY, bold=True)
    r2 = p.add_run(); set_run(r2, blind, 12, INK_VARIANT, FONT_BODY)
    p = tf.add_paragraph(); p.space_after = Pt(0)
    r3 = p.add_run(); set_run(r3, "Preuve à mener : ", 12, TEAL, FONT_BODY, bold=True)
    r4 = p.add_run(); set_run(r4, exp, 12, TEAL, FONT_BODY)
pagefoot(s, PAGE)
notes(s, """
Trois premiers risques, chacun avec son angle aveugle et l'EXPÉRIENCE qui le tranche (kill-criteria) :

1. UNIT ECONOMICS — Le risque numéro un. Chaque diagnostic gratuit consomme du LLM (intake, narration, assistant) et du Firecrawl (feed prix, veille catalogue). Si le taux de conversion est bas, on paie de l'IA pour des leads qui ne convertissent pas. Preuve : instrumenter le coût IA réel par diagnostic sur 50 parcours, croiser avec le taux de conversion P1, et tester l'effet de quotas + cache + modèle économe. Kill-criteria : si coût IA/lead non amorti même avec quotas et conversion réaliste → revoir le free tier.

2. DÉPENDANCE FOURNISSEURS — Le paradoxe : on vend la souveraineté mais le cerveau (LLM) et les yeux (Firecrawl) sont des tiers. Atténuations DÉJÀ dans le code : (a) tous les appels LLM ont un repli déterministe (le produit fonctionne si le LLM est KO — c'est testé) ; (b) le calcul de coût est 100 % déterministe (le LLM ne calcule jamais) ; (c) LiteLLM permet de changer de provider, y compris vers un modèle auto-hébergé, sans toucher l'app. Preuve : faire tourner un LLM open auto-hébergé via LiteLLM et mesurer la qualité/coût ; documenter la dégradation gracieuse comme argument de vente souverain.

3. ÉROSION DES MOATS — (a) Exit Escrow copiable : oui techniquement, mais un acteur dont le modèle EST le lock-in ne le fera pas (anti-incitatif). (b) Fiduciaire : un commissionné ne peut pas copier sans saborder son revenu (vrai moat défensif). (c) Effet réseau : c'est le maillon faible — il n'existe PAS tant qu'il n'y a pas un volume de déploiements monitorés (Lot 2). Preuve : estimer le nombre de déploiements nécessaires pour que la calibration batte le ±30 % ; jusque-là, ne pas survendre le Network comme un moat actif.
""")

# --- 18. AVOCAT DU DIABLE — risques détaillés (table 2) --------------------
PAGE += 1
s = add_slide()
std_header(s, "Avocat du diable · risques 2/2", "Risque → angle aveugle → preuve à mener", accent=ERROR, eb_color=ERROR)
risk_rows2 = [
    ("Risque juridique", "Conseil souveraineté/conformité peut être lu comme du conseil juridique ou financier réglementé.", "Disclaimer « ingénierie, pas juridique » partout ; partenariat DPO/avocat ; faire relire les CGU."),
    ("Marché & maturité", "Aucun TAM sourcé : « souveraineté IA » peut être un marché de niche, ou pas encore payant.", "Sourcer une vraie taille de marché ; valider 5 LOI/arrhes P1 avant d'extrapoler."),
    ("Cycle B2B & exécution solo", "Vente régulés = cycle long ; petite équipe = risque d'exécution et de bus-factor.", "Concentrer P1/P2 (cycle court) ; documenter ; recruter sur les goulots ; jalonner."),
    ("Le piège du [PLACEHOLDER]", "Si la WTP < coût de servir la cuisine, le modèle payant ne tient pas.", "Clore le sondage (seuils VW/conjoint) AVANT tout scale ; tester 3 prix sur ventes réelles."),
]
top = Inches(2.0); rh = Inches(1.12); gap = Inches(0.12); l = Inches(0.85); w = Inches(11.6)
for i, (name, blind, exp) in enumerate(risk_rows2):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    c = card(s, l, t, w, rh)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.09), rh)
    bar.fill.solid(); bar.fill.fore_color.rgb = ERROR; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    para(tf, name, 14, ERROR, FONT_TITLE, bold=True, first=True, space_after=3)
    p = tf.add_paragraph(); p.space_after = Pt(3)
    r1 = p.add_run(); set_run(r1, "Angle aveugle : ", 11.5, INK, FONT_BODY, bold=True)
    r2 = p.add_run(); set_run(r2, blind, 11.5, INK_VARIANT, FONT_BODY)
    p = tf.add_paragraph(); p.space_after = Pt(0)
    r3 = p.add_run(); set_run(r3, "Preuve : ", 11.5, TEAL, FONT_BODY, bold=True)
    r4 = p.add_run(); set_run(r4, exp, 11.5, TEAL, FONT_BODY)
pagefoot(s, PAGE)
notes(s, """
Quatre risques restants :

4. RISQUE JURIDIQUE — Strate conseille des choix de souveraineté et de conformité (RGPD/CNDP). Si le discours glisse vers « voici votre obligation légale », on frôle le conseil juridique réglementé. De même, parler d'« économie » et de « négo vendor » peut frôler le conseil financier. Mitigations : disclaimer « ingénierie, pas juridique » présent dans la roadmap résidence/DR ; partenariat DPO/avocat (slide partenariats) ; faire relire les CGU et la charte fiduciaire par un juriste. Kill-criteria : si un conseil produit pouvait être requalifié en acte réglementé → restreindre le périmètre.

5. MARCHÉ & MATURITÉ — Honnêteté brutale : on n'a PAS de TAM sourcé pour « souveraineté IA des bases mémorielles ». Le marché peut être réel mais de niche, ou pas encore prêt à payer. On ne doit pas faire semblant. Preuve : (a) sourcer une vraie estimation de marché (rapports analystes, sans l'inventer) ; (b) surtout, valider la demande par 5 LOI/arrhes réelles sur P1 AVANT d'extrapoler quoi que ce soit. La validation terrain prime sur le TAM théorique.

6. CYCLE B2B & EXÉCUTION SOLO — Les régulés (P4) ont un cycle long (DPA, audits) ; commencer par là tuerait le runway. D'où la focalisation P1/P2 (cycle court). Risque d'exécution : petite équipe = bus-factor élevé. Mitigations : documentation (ADR, progress.md déjà tenus), jalonnement, recrutement ciblé sur les goulots (support, vente).

7. LE PIÈGE DU [PLACEHOLDER] — LE risque transversal. Si la willingness-to-pay pour la « cuisine » est inférieure au coût de servir (IA + support + déploiement), le modèle payant ne tient pas, même avec un bon produit. C'est pourquoi le sondage WTP n'est pas un confort mais un GATE : clore le sondage (≥ 15 VW + ≥ 30 conjoint) avant tout scale, puis tester 3 prix sur des ventes réelles. Kill-criteria : si la WTP médiane < coût de servir → repenser le packaging (ex. plus de self-service, moins de service humain).

Autres angles à mentionner si on a le temps : cannibalisation conseil vs produit (le livrable gratuit peut suffire à certains → assumé par « recette ouverte ») ; concentration géographique (Maroc-UE) ; scénarios de sortie/financement (bootstrap vs levée — à trancher selon le sondage et les 5 ventes).
""")

# --- 19. QUESTIONS OUVERTES & KILL-CRITERIA --------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Décisions à trancher & kill-criteria", "Les questions ouvertes (l'utilisateur tranche)", accent=BLUE, eb_color=BLUE)
qs = [
    ("Prix de la cuisine", "Quel forfait Pro ? → après sondage (VW + conjoint, seuils atteints).", BLUE),
    ("Bootstrap ou levée", "Financer le Lot 2 sur fonds propres / revenus, ou lever ?", GOLD_TXT),
    ("Free tier IA", "Quotas / cache / modèle économe pour borner le coût IA des gratuits ?", TEAL),
    ("Ordre des segments", "P3 (agences) ou P4 (régulés) en second, après P1/P2 ?", BLUE),
    ("Périmètre fiduciaire", "Rémunération de négo vendor : part de l'économie, ou forfait ?", GOLD_TXT),
    ("Effet réseau", "Combien de déploiements pour que la calibration batte le ±30 % ?", TEAL),
]
cw, ch, gap, top, l0 = Inches(3.7), Inches(1.7), Inches(0.25), Inches(2.05), Inches(0.85)
for i, (h, body, col) in enumerate(qs):
    row, coln = divmod(i, 3)
    l = Emu(int(l0) + coln * (int(cw) + int(gap)))
    t = Emu(int(top) + row * (int(ch) + int(Inches(0.22))))
    c = card(s, l, t, cw, ch, fill=SURFACE_LOW)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, cw, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.22)
    para(tf, h, 14.5, INK, FONT_TITLE, bold=True, first=True, space_after=5)
    para(tf, body, 12, INK_VARIANT, space_after=0)
pagefoot(s, PAGE)
notes(s, """
Les questions ouvertes structurantes, présentées comme des DÉCISIONS À TRANCHER par Amine (principe « avis non orienté » : on expose les options et les preuves nécessaires, on ne tranche pas à sa place).

1. Prix de la cuisine (Pro) : décision GATÉE par le sondage. Ne rien fixer avant les seuils (≥ 15 VW, ≥ 30 conjoint).
2. Bootstrap ou levée : le Lot 2 (déploiement assisté + monitoring) demande de l'effort. Le financer sur revenus (plus lent, plus souverain — cohérent avec le récit) ou lever (plus rapide, dilution + pression de croissance qui peut heurter le positionnement fiduciaire) ? À trancher selon le sondage et les 5 ventes.
3. Free tier IA : quotas, cache, modèle économe — décision produit qui pèse directement sur les unit economics.
4. Ordre des segments : P3 (agences, multi-tenancy déjà câblé) ou P4 (régulés, premium mais cycle long) en second ?
5. Périmètre fiduciaire de la négo vendor : facturer une part de l'économie générée (aligné valeur, à divulguer) ou un forfait (plus simple, moins de risque de conflit perçu) ?
6. Effet réseau : combien de déploiements monitorés pour que la calibration collective batte vraiment le ±30 % ? Tant qu'on ne sait pas, ne pas survendre le moat ③.

KILL-CRITERIA récapitulés (les seuils d'abandon/pivot) :
- WTP médiane < coût de servir la cuisine → repenser le packaging.
- Coût IA/lead non amorti même avec quotas + conversion réaliste → revoir le free tier.
- < 5 LOI/arrhes P1 après présentation → le wedge n'est pas validé, repenser la cible.
- Un conseil produit requalifiable en acte réglementé → restreindre le périmètre.

Le rôle de Strate (et de ce deck) est de présenter l'analyse neutre ; Amine décide.
""")

# --- 20. ROADMAP & ASKS -----------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Roadmap & prochaines étapes", "Jalons de preuve, pas seulement de features", accent=TEAL)
# timeline 3 lots + état
lots = [
    ("Lot 1 — Conseil + moats", "LIVRÉ + déployé en prod. Reco 7 couches sourcée, ensemble, livrable, Exit Escrow, charte fiduciaire, rails Network, épic LLM complet (intake/narration/assistant/reco vivante/admin).", TEAL, "FAIT"),
    ("Lot 2 — Cuisine payante", "Provisioning hybride (human-in-the-loop), coffre de secrets, monitoring (Infra Health Score + MEL), Network actif. C'est ici que naissent les revenus récurrents.", BLUE, "À VENIR"),
    ("En cours / proche", "Résidence & continuité (DR multi-région), conformité transferts (RGPD chap. V / CNDP / Cloud Act) — épic S-043→048.", GOLD_TXT, "PLANIFIÉ"),
]
top = Inches(2.0); rh = Inches(1.25); gap = Inches(0.16); l = Inches(0.85); w = Inches(11.6)
for i, (name, body, col, tag) in enumerate(lots):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    c = card(s, l, t, w, rh)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.09), rh)
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(2.3)
    para(tf, name, 15, INK, FONT_TITLE, bold=True, first=True, space_after=4)
    para(tf, body, 12, INK_VARIANT, space_after=0)
    tb2, tf2 = textbox(s, Emu(int(l) + int(w) - int(Inches(2.1))), t, Inches(1.9), rh, anchor=MSO_ANCHOR.MIDDLE)
    pp = tf2.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run(); set_run(rr, tag, 12, col, FONT_MONO, bold=True)
# asks
tb, tf = textbox(s, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.7))
para(tf, "Asks : (1) clore le sondage WTP · (2) 5 LOI/arrhes P1 · (3) arbitrer bootstrap/levée pour financer le Lot 2.",
     14, TEAL, FONT_BODY, bold=True, first=True)
pagefoot(s, PAGE)
notes(s, """
La roadmap est tenue par les JALONS DE PREUVE, pas seulement par les features.

Lot 1 (LIVRÉ + déployé en prod, infra.ai-mpower.com) : tout le « conseil + moats ». Reco déterministe 7 couches sourcée ±30 %, ensemble multi-config, livrable MD/PDF, Exit Escrow (testé), charte fiduciaire, rails du Network (RLS, consentement). L'épic LLM est complet : intake libre → Profile (validé serveur), narration (chiffres préservés), assistant Q&A contextuel, reco « vivante » (catalogue sourcé web + LLM, garde-fou DÉFCON 1), console admin de prompts versionnés. C'est VENDABLE SEUL.

Lot 2 (à venir) : la « cuisine payante » et donc le revenu récurrent. Provisioning hybride human-in-the-loop (l'utilisateur crée le compte + saisit sa carte, l'agent optimise APRÈS via OAuth/MCP — jamais d'inscription ni de carte par l'agent → CGU + responsabilité maîtrisées), coffre de secrets, monitoring (Infra Health Score type NEWS2 + Dispatch Deviation Guide type MEL aviation), Network actif (collecte du coût réel → recalibration). C'est l'investissement clé à financer.

En cours / proche : épic résidence & continuité (DR multi-région, table de conformité des transferts UE/Maroc/US), avec disclaimer « ingénierie, pas juridique ».

ASKS (ce dont Amine a besoin pour avancer) :
1. Clore le sondage WTP (seuils VW + conjoint) → débloque le prix et le P&L réel.
2. Obtenir 5 LOI/arrhes sur P1 → valide la demande et l'offre.
3. Arbitrer bootstrap vs levée → finance le Lot 2 (le moteur de revenu récurrent).

Jalons de preuve = ces 3 asks + la première validation de l'Exit Escrow sur une stack réelle (banc d'essai : la stack d'Amine).
""")

# --- 21. SYNTHÈSE DÉCISIONS -------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Recommandations structurantes", "5 décisions à trancher (options, pas verdicts)", accent=TEAL)
decs = [
    ("Faire du sondage WTP un GATE", "Ne fixer aucun prix ni P&L avant les seuils. Le [PLACEHOLDER] reste assumé."),
    ("Tenir le wedge P1 d'abord", "Refuser P5, séquencer P3/P4 après ; CAC quasi nul = runway préservé."),
    ("Borner le coût IA du free tier", "Quotas / cache / modèle économe via LiteLLM avant tout scale d'acquisition."),
    ("Faire de l'Exit Escrow le héros", "Centrer la comm sur la seule preuve DURE ; démonstration live = moment wow."),
    ("Financer le Lot 2 sur jalons", "Bootstrap si possible (cohérent avec le récit souverain) ; lever seulement si la WTP + 5 ventes le justifient."),
]
top = Inches(2.05); rh = Inches(0.84); gap = Inches(0.13); l = Inches(0.85); w = Inches(11.6)
for i, (h, body) in enumerate(decs):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    c = card(s, l, t, w, rh)
    # numéro pastille
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(l) + int(Inches(0.18))), Emu(int(t) + int(Inches(0.2))), Inches(0.44), Inches(0.44))
    num.fill.solid(); num.fill.fore_color.rgb = TEAL; num.line.fill.background(); num.shadow.inherit = False
    ntf = num.text_frame; _no_autofit(ntf); ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    npp = ntf.paragraphs[0]; npp.alignment = PP_ALIGN.CENTER
    nr = npp.add_run(); set_run(nr, str(i + 1), 16, WHITE, FONT_TITLE, bold=True)
    tf = c.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.85); tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]
    r1 = p.add_run(); set_run(r1, h + " — ", 14.5, INK, FONT_TITLE, bold=True)
    r2 = p.add_run(); set_run(r2, body, 12.5, INK_VARIANT, FONT_BODY)
pagefoot(s, PAGE)
notes(s, """
Les 5 décisions structurantes, présentées comme des OPTIONS recommandées à trancher (jamais un verdict imposé — principe d'Amine).

1. Sondage WTP = GATE absolu. La discipline DÉFCON 1 appliquée au business : pas de prix au doigt mouillé. Tant que le sondage n'a pas ses seuils, le prix reste [PLACEHOLDER] et le P&L reste paramétrique. C'est inconfortable mais honnête, et ça évite de bâtir sur du sable.

2. Tenir le wedge P1. La tentation sera d'aller chercher les gros (P4/P5) qui « valent plus ». Mais le cycle long et la conformité lourde tueraient le runway. P1 (communauté) = demande quasi acquise, CAC ~0, preuve sociale. Séquencer le reste.

3. Borner le coût IA du free tier. C'est l'angle aveugle économique principal. Décider tôt d'une politique de quotas/cache/modèle économe (LiteLLM facilite le choix du modèle côté proxy) pour ne pas se faire saigner par les gratuits.

4. Faire de l'Exit Escrow le héros de la communication. C'est la seule preuve DURE. Une démonstration live (générer un bundle, le redéployer ailleurs) est le « moment wow » qui matérialise l'anti-lock-in. Tout le reste = « à mesurer chez vous ».

5. Financer le Lot 2 sur jalons. Bootstrap si la trésorerie le permet (cohérent avec le récit souverain et le positionnement fiduciaire — pas de pression d'investisseur qui pousserait à trahir la charte). Lever seulement si la WTP validée + les 5 ventes prouvent un modèle scalable qui justifie l'accélération.

Insister : ce sont des recommandations argumentées, pas des ordres. Amine arbitre avec les preuves en main.
""")

# --- 22. ANNEXE — hypothèses du P&L ----------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Annexe · 1/3", "Hypothèses du P&L (à valider, non engageantes)", accent=CARD_BORDER, eb_color=INK_VARIANT)
tb, tf = textbox(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(4.7))
hyps = [
    ("ARPA", "= prix service Pro /mois = [PLACEHOLDER]. Souverain+ sur devis (ARPA supérieur, cycle plus long). Non instancié tant que le sondage n'est pas clos."),
    ("Clients payants (scénarios)", "Prudent 8/25/60 · Base 15/60/180 · Ambitieux 30/150/500 sur 3 ans. Hypothèses de modélisation, PAS des prévisions."),
    ("Marge brute", "Service quasi pur (infra reversée sans marge). Élevée en théorie, mais nette des coûts variables IA (LLM + Firecrawl)."),
    ("Conversion / churn", "Cibles HYPOTHÉTIQUES : conversion P1 > 10 %, rétention > 80 % (comptes avec Exit bundle). À observer sur ventes réelles."),
    ("CAC", "≈ 0 sur P1 (communauté) ; non modélisé hors réseau — à mesurer."),
    ("Coûts variables IA", "Intake + narration + assistant + veille (LLM via LiteLLM) + price feed (Firecrawl). Dette transparente, non chiffrée dans le moteur (story future)."),
    ("OPEX", "Dev (équipe réduite), support, hébergement plateforme (self-host Coolify/serveur). Léger au démarrage."),
]
for i, (k, v) in enumerate(hyps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8); p.line_spacing = 1.05
    r1 = p.add_run(); set_run(r1, k + " : ", 13, TEAL, FONT_TITLE, bold=True)
    r2 = p.add_run(); set_run(r2, v, 12.5, INK_VARIANT, FONT_BODY)
pagefoot(s, PAGE)
notes(s, """
Le détail des hypothèses du P&L, pour transparence totale. À répéter : aucune n'est engageante ; toutes sont à valider.

- ARPA : c'est le [PLACEHOLDER]. Le sondage Van Westendorp + conjoint le calera. Souverain+ aura un ARPA plus élevé mais un cycle de vente plus long (devis).
- Clients payants : les 3 trajectoires (prudent/base/ambitieux) sont des structures de scénario, choisies pour borner le raisonnement, pas des projections de marché.
- Marge brute : structurellement bonne (le revenu Strate est du service, pas de la revente d'infra), MAIS il faut la calculer NETTE des coûts variables IA. C'est la nuance qui distingue la marge théorique de la marge réelle.
- Conversion / churn : cibles du PRD §10, à traiter comme hypothèses à observer.
- CAC : ~0 sur P1 ; inconnu hors réseau.
- Coûts variables IA : explicitement une dette transparente du projet (le coût LLM à l'usage est signalé non chiffré dans le moteur). À instrumenter.
- OPEX : léger (équipe réduite, hébergement maîtrisé), ce qui allonge le runway en bootstrap.

Tout chiffre absolu en euros est volontairement absent : il dépend de l'ARPA inconnu.
""")

# --- 23. ANNEXE — glossaire -------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Annexe · 2/3", "Glossaire", accent=CARD_BORDER, eb_color=INK_VARIANT)
gloss = [
    ("Base mémorielle", "Mémoire d'organisation IA : doctrine, historique, savoir — souveraine et persistante."),
    ("7 couches (C0→C6)", "Contrat YAML, surface MCP, orchestrateur RAG, retrieval/rerank, embeddings, stockage bi-temporel, infra."),
    ("Preset", "Palier d'architecture (LIGHT / MEDIUM / HARD) dérivé d'un score de besoin explicable."),
    ("Exit Escrow", "Bundle reproductible (IaC + dumps + vault + runbook) redéployable ailleurs — moat ①."),
    ("Fiduciary Mode", "Engagement de mandataire : zéro commission vendor cachée — moat ②."),
    ("Intelligence Network", "Effet réseau de coûts réels anonymisés (opt-in RGPD) — moat ③."),
    ("Ensemble multi-config", "N variantes (souveraineté max / coût min / time-to-V1) ; le spread = l'incertitude."),
    ("Van Westendorp / CBC", "Méthodes de mesure du prix psychologique (VW) et des arbitrages de valeur (conjoint)."),
    ("±30 %", "Bande d'incertitude assumée du chiffrage ; recalibrée par le Network ; jamais une prédiction de conso."),
    ("DÉFCON 1", "Règle interne : zéro donnée erronée présentée comme fiable ; tout chiffre sourcé ou étiqueté hypothèse."),
]
# deux colonnes
tb, tf = textbox(s, Inches(0.85), Inches(2.0), Inches(5.7), Inches(4.7))
tb2, tf2 = textbox(s, Inches(6.75), Inches(2.0), Inches(5.7), Inches(4.7))
half = (len(gloss) + 1) // 2
for i, (k, v) in enumerate(gloss):
    target = tf if i < half else tf2
    first = (i == 0) or (i == half)
    p = target.paragraphs[0] if first else target.add_paragraph()
    p.space_after = Pt(9); p.line_spacing = 1.05
    r1 = p.add_run(); set_run(r1, k, 13, TEAL, FONT_TITLE, bold=True)
    r2 = p.add_run(); set_run(r2, " — " + v, 11.5, INK_VARIANT, FONT_BODY)
pagefoot(s, PAGE)
notes(s, """
Glossaire pour aligner le vocabulaire avec un public non-technique (investisseurs, partenaires). Les termes clés :
- Base mémorielle = le « produit » du client (sa mémoire d'organisation), pas l'outil Strate.
- 7 couches = l'architecture que Strate recommande et chiffre.
- Preset = le palier (LIGHT/MEDIUM/HARD), désormais dérivé d'un score explicable (pas une boîte noire — cf. ADR-020).
- Les 3 moats sont rappelés.
- Van Westendorp / CBC = les méthodes du sondage de prix.
- ±30 % = l'honnêteté assumée du chiffrage (c'est un simulateur, pas une boule de cristal).
- DÉFCON 1 = la règle de rigueur qui traverse tout le projet ET ce deck.
Utiliser cette slide en référence si une question de vocabulaire surgit.
""")

# --- 24. ANNEXE — sources ---------------------------------------------------
PAGE += 1
s = add_slide()
std_header(s, "Annexe · 3/3", "Sources & provenance", accent=CARD_BORDER, eb_color=INK_VARIANT)
tb, tf = textbox(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(4.7))
para(tf, "Sources internes (dépôt Strate)", 14, INK, FONT_TITLE, bold=True, first=True, space_after=7)
internal = [
    "PRD.md — spec, JTBD, personas, périmètre, modèle éco, métriques, risques.",
    "docs/MOAT-HUNT.md — les 3 moats, analogies inter-industries, scores.",
    "docs/DECISIONS.md — ADR-001→020 (presets scorés, déterminisme, etc.).",
    "docs/pricing/wtp-research.md — méthodo Van Westendorp + conjoint, paliers candidats.",
    ".ralph/prd.json + progress.md — état réel des features (épic LLM, backup, compute).",
    "lib/fiduciary/charter.ts — charte fiduciaire (zéro commission), testée.",
    "design-reference/mn_mo_brand_identity/DESIGN.md — identité de marque (palette, typo).",
]
for it in internal:
    bullet(tf, it, size=12, space_after=5, marker_color=TEAL)
para(tf, "Sources externes (analogies de moats, dans le PRD / MOAT-HUNT)", 14, INK, FONT_TITLE, bold=True, space_after=7, space_before=8)
external = [
    "Escrow logiciel / IP Bankruptcy Act — en.wikipedia.org/wiki/Source_code_escrow",
    "Courtier & secret commission — brabners.com ; fosterec.com ; moneytothemasses (Flipper).",
    "ISAC / collective defense — anomali.com ; nationalisacs.org.",
    "Ensemble forecasting — science.org (FuXi-ENS) · MEL aviation — skybrary.aero · NEWS2 — nice.org.uk.",
]
for it in external:
    bullet(tf, it, size=12, space_after=5, marker_color=BLUE)
para(tf, "Aucune taille de marché (TAM) n'est sourcée → traitée comme inconnue, pas comme acquis (cf. risques).",
     11.5, ERROR, FONT_BODY, italic=True, space_before=8)
pagefoot(s, PAGE)
notes(s, """
Traçabilité complète (exigence DÉFCON 1 : citer les sources).

Sources INTERNES : tout le business model est dérivé du dépôt réel — PRD, MOAT-HUNT, DECISIONS (ADR), wtp-research, prd.json/progress.md (état des features), charter.ts (charte fiduciaire), DESIGN.md (identité). Aucune feature inventée.

Sources EXTERNES : ce sont les analogies qui fondent les moats (escrow logiciel, courtier fiduciaire/leçon Flipper, ISAC, ensemble forecasting, MEL, NEWS2). Elles sont listées et liées dans le PRD et MOAT-HUNT. Ce sont des MODÈLES MENTAUX validés sur d'autres industries, pas des garanties de succès.

POINT D'HONNÊTETÉ FINAL : aucune taille de marché (TAM/SAM/SOM) n'est sourcée dans ce deck. C'est délibéré — on ne fabrique pas de stat de marché. La taille du marché « souveraineté IA » est traitée comme une INCONNUE à sourcer/valider (slide risques), pas comme un acquis. C'est cohérent avec la règle « zéro donnée erronée présentée comme fiable ».
""")

# --- 25. CLÔTURE ------------------------------------------------------------
PAGE += 1
s = add_slide(INVERSE)
accent_bar(s, TEAL)
tb, tf = textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Une thèse claire, des moats défendables,", 30, WHITE, FONT_TITLE, bold=True, first=True, space_after=2)
para(tf, "des chiffres honnêtement étiquetés.", 30, RGBColor(0x6B, 0xD8, 0xCB), FONT_TITLE, bold=True, space_after=14)
para(tf, "Prochain pas : clore le sondage, signer 5 clients P1, financer la cuisine.", 16, INVERSE_TXT, FONT_BODY)
tb, tf = textbox(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
para(tf, "Strate · base mémorielle souveraine · sans migration, sans verrouillage, sans coûts cachés.", 12, RGBColor(0xB4, 0xC5, 0xFF), FONT_MONO, first=True)
notes(s, """
Clôture. Récapituler en 3 temps :
1. La thèse est claire : le Cloudflare de la mémoire IA souveraine — anti-lock-in vendu comme un produit.
2. Les moats sont défendables : Exit Escrow (preuve dure), Fiduciary (barrière que les commissionnés ne peuvent copier), Network (effet réseau, à activer au Lot 2).
3. Les chiffres sont honnêtement étiquetés : prix = [PLACEHOLDER], P&L = hypothèses, marché = inconnue à valider. Cette honnêteté est un actif de crédibilité, pas une faiblesse.

Appel à l'action / prochains pas : clore le sondage WTP, obtenir 5 clients/LOI P1, financer le Lot 2 (la cuisine = le revenu récurrent). C'est la séquence qui transforme un produit livré en entreprise viable.

Fin de la présentation.
""")
pagefoot(s, PAGE, "")

# ---------------------------------------------------------------------------
import os
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "strate-business-model.pptx")
prs.save(OUT)
print("OK ->", OUT)
print("slides:", len(prs.slides._sldIdLst))
