# -*- coding: utf-8 -*-
"""Vérifie le .pptx généré : nombre de slides, présence de notes, titres."""
import os
from pptx import Presentation

PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "strate-business-model.pptx")
prs = Presentation(PATH)

n = len(prs.slides)
print(f"Fichier : {PATH}")
print(f"Slides  : {n}")
print(f"Taille  : {os.path.getsize(PATH)} octets")
print("-" * 70)

missing_notes = []
empty_notes = []
for i, slide in enumerate(prs.slides, start=1):
    # premier texte non vide = titre approximatif
    title_txt = ""
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip():
            title_txt = shp.text_frame.text.strip().split("\n")[0]
            break
    has_notes = slide.has_notes_slide
    note_len = 0
    if has_notes:
        note_len = len(slide.notes_slide.notes_text_frame.text.strip())
    else:
        missing_notes.append(i)
    if has_notes and note_len < 40:
        empty_notes.append((i, note_len))
    flag = "OK " if (has_notes and note_len >= 40) else "!! "
    print(f"{flag}[{str(i).zfill(2)}] notes={note_len:>5} | {title_txt[:58]}")

print("-" * 70)
print(f"Slides SANS notes        : {missing_notes if missing_notes else 'aucune'}")
print(f"Slides notes trop courtes: {empty_notes if empty_notes else 'aucune'}")
all_ok = (not missing_notes) and (not empty_notes)
print(f"RESULTAT : {'TOUTES les slides ont des notes substantielles' if all_ok else 'PROBLEME detecte'}")
