# -*- coding: utf-8 -*-
import os
from pptx import Presentation
PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "strate-business-model-canvas.pptx")
prs = Presentation(PATH)
print("Fichier :", PATH)
print("Slides  :", len(prs.slides))
print("Taille  :", os.path.getsize(PATH), "octets")
print("-" * 64)
missing = []
for i, slide in enumerate(prs.slides, 1):
    title = ""
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip():
            title = shp.text_frame.text.strip().split("\n")[0]
            break
    has = slide.has_notes_slide
    nlen = len(slide.notes_slide.notes_text_frame.text.strip()) if has else 0
    if not has or nlen < 40:
        missing.append(i)
    flag = "OK " if (has and nlen >= 40) else "!! "
    print(f"{flag}[{str(i).zfill(2)}] notes={nlen:>5} | {title[:54]}")
print("-" * 64)
print("RESULTAT :", "toutes les slides ont des notes substantielles" if not missing else f"PROBLEME slides {missing}")
