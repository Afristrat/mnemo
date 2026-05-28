# -*- coding: utf-8 -*-
"""
Génère le Business Model Canvas (Osterwalder) de Strate, rempli.
Sortie : docs/strate-business-model-canvas.pptx

Disposition BMC standard à 9 blocs :
  Rangée haute : Partenaires clés | Activités clés / Ressources clés | Proposition de valeur | Relations clients / Canaux | Segments de clientèle
  Rangée basse pleine largeur : Structure de coûts | Flux de revenus

Palette/typo de design-reference/mn_mo_brand_identity/DESIGN.md.
Garde-fous DÉFCON 1 : prix = [PLACEHOLDER] ; tout revenu/coût étiqueté hypothèse ;
aucune stat inventée ; produit fidèle au dépôt (Lot 1 livré + épic LLM ; Lot 2/résidence à venir).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Palette (DESIGN.md)
TEAL        = RGBColor(0x00, 0x68, 0x5F)
TEAL_DIM    = RGBColor(0x00, 0x83, 0x78)
BLUE        = RGBColor(0x00, 0x51, 0xD5)
GOLD        = RGBColor(0x9D, 0x6A, 0x00)
GOLD_TXT    = RGBColor(0x7D, 0x54, 0x00)
INK         = RGBColor(0x13, 0x1B, 0x2E)
INK_VARIANT = RGBColor(0x3D, 0x49, 0x47)
SURFACE     = RGBColor(0xFA, 0xF8, 0xFF)
SURFACE_LOW = RGBColor(0xF2, 0xF3, 0xFF)
CARD_BORDER = RGBColor(0xBC, 0xC9, 0xC6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ERROR       = RGBColor(0xBA, 0x1A, 0x1A)
INVERSE     = RGBColor(0x28, 0x30, 0x44)
INVERSE_TXT = RGBColor(0xEE, 0xF0, 0xFF)
MINT        = RGBColor(0x6B, 0xD8, 0xCB)

FONT_TITLE = "Space Grotesk"
FONT_BODY  = "Inter"
FONT_MONO  = "JetBrains Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_slide(bg=SURFACE):
    s = prs.slides.add_slide(BLANK)
    _bg(s, bg)
    return s


def set_run(run, text, size, color, font=FONT_BODY, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, color, font=FONT_BODY, bold=False, italic=False,
         sa=4, sb=0, align=PP_ALIGN.LEFT, first=False, ls=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(sa); p.space_before = Pt(sb)
    if ls:
        p.line_spacing = ls
    r = p.add_run()
    set_run(r, text, size, color, font, bold, italic)
    return p


def rrect(slide, l, t, w, h, fill, border=None, border_w=1.0, radius=0.05):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if border is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = border; c.line.width = Pt(border_w)
    c.shadow.inherit = False
    try:
        c.adjustments[0] = radius
    except Exception:
        pass
    return c


def rect(slide, l, t, w, h, fill, border=None, border_w=1.0):
    c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if border is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = border; c.line.width = Pt(border_w)
    c.shadow.inherit = False
    return c


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def bmc_block(slide, l, t, w, h, num, heading, head_color, items, dense=False,
              head_fs=11.5, item_fs=9.0, sa=2.4):
    """Un bloc du canevas : carte blanche, bandeau d'en-tête coloré, puces."""
    rrect(slide, l, t, w, h, WHITE, border=CARD_BORDER, border_w=1.0, radius=0.035)
    # bandeau d'en-tête
    head_h = Inches(0.34)
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, head_h)
    hd.fill.solid(); hd.fill.fore_color.rgb = head_color; hd.line.fill.background(); hd.shadow.inherit = False
    htf = hd.text_frame; htf.word_wrap = True; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    htf.margin_left = Inches(0.1); htf.margin_right = Inches(0.06); htf.margin_top = 0; htf.margin_bottom = 0
    hp = htf.paragraphs[0]
    r1 = hp.add_run(); set_run(r1, num + "  ", 11.5, WHITE, FONT_MONO, bold=True)
    r2 = hp.add_run(); set_run(r2, heading, head_fs, WHITE, FONT_TITLE, bold=True)
    # corps
    body_t = Emu(int(t) + int(head_h) + int(Inches(0.06)))
    tb, tf = textbox(slide, Emu(int(l) + int(Inches(0.1))), body_t,
                     Emu(int(w) - int(Inches(0.2))), Emu(int(h) - int(head_h) - int(Inches(0.12))))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(sa); p.line_spacing = 1.0
        rm = p.add_run(); set_run(rm, "› ", item_fs, head_color, FONT_BODY, bold=True)
        # support du gras en tête via tuple (lead, rest)
        if isinstance(it, tuple):
            rb = p.add_run(); set_run(rb, it[0], item_fs, INK, FONT_BODY, bold=True)
            rt = p.add_run(); set_run(rt, it[1], item_fs, INK_VARIANT, FONT_BODY)
        else:
            rt = p.add_run(); set_run(rt, it, item_fs, INK_VARIANT, FONT_BODY)


# ===========================================================================
# SLIDE 1 — TITRE
# ===========================================================================
s = add_slide(INVERSE)
bar = rect(s, Inches(0), Inches(0), Inches(0.16), Inches(7.5), TEAL)
tb, tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.4))
p = para(tf, "BUSINESS MODEL CANVAS · OSTERWALDER · DOCUMENT DE TRAVAIL 2026", 13, MINT, FONT_BODY, bold=True, first=True)
for r in p.runs:
    r._r.get_or_add_rPr().set('spc', '160')
tb, tf = textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.5))
para(tf, "Strate", 60, WHITE, FONT_TITLE, bold=True, first=True, sa=2)
para(tf, "Le canevas en 9 blocs d'une base mémorielle souveraine.", 21, INVERSE_TXT, FONT_TITLE)
tb, tf = textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.6))
para(tf, "Configurateur SaaS souverain d'infrastructure de « base mémorielle » IA · conseil → déploiement · recette ouverte, cuisine payante.", 15, INVERSE_TXT, FONT_BODY, first=True, sa=6)
para(tf, "Trois moats (playbook Cloudflare) : ① Exit Escrow · ② Fiduciary Mode · ③ Intelligence Network.", 14, RGBColor(0xB4, 0xC5, 0xFF), FONT_BODY)
tb, tf = textbox(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4))
para(tf, "Prix = [PLACEHOLDER] (sondage en cours) · tout revenu/coût = HYPOTHÈSE à valider · aucune stat inventée.", 11, MINT, FONT_MONO, first=True)
notes(s, """
Slide d'ouverture du Business Model Canvas (méthode Osterwalder & Pigneur, Business Model Generation).

Le canevas qui suit synthétise sur UNE page les 9 briques du modèle économique de Strate. Il est dérivé du produit RÉEL (dépôt Infra) : Lot 1 (conseil + moats) livré et déployé en prod, épic LLM complet ; Lot 2 (cuisine payante : déploiement assisté + monitoring) et l'épic résidence/DR sont signalés « à venir ».

Garde-fous (règle DÉFCON 1 d'Amine) à rappeler d'entrée :
- Le prix de vente du service Strate est un [PLACEHOLDER] : un sondage de willingness-to-pay (Van Westendorp + conjoint) est en cours (cf. docs/pricing/wtp-research.md). Aucun montant ferme n'est asséné.
- Toute mention de revenus ou de coûts est une HYPOTHÈSE de modélisation, pas une prévision.
- Aucune taille de marché n'est inventée.

Les slides 3 à 5 détaillent chaque bloc, avec le « pourquoi » et les preuves. La slide 2 est le canevas grille à présenter en plénière.
""")

# ===========================================================================
# SLIDE 2 — LE CANEVAS (grille 9 blocs)
# ===========================================================================
s = add_slide(SURFACE)
# en-tête mince
tb, tf = textbox(s, Inches(0.35), Inches(0.18), Inches(12.6), Inches(0.5))
p = tf.paragraphs[0]
r1 = p.add_run(); set_run(r1, "Business Model Canvas — ", 19, INK, FONT_TITLE, bold=True)
r2 = p.add_run(); set_run(r2, "Strate", 19, TEAL, FONT_TITLE, bold=True)
tb, tf = textbox(s, Inches(10.0), Inches(0.24), Inches(2.95), Inches(0.4))
para(tf, "prix = [PLACEHOLDER]", 10, ERROR, FONT_MONO, first=True, align=PP_ALIGN.RIGHT)

# Géométrie de la grille
GX = Inches(0.35)          # marge gauche
GY = Inches(0.78)          # haut de la grille
GW = Inches(12.63)         # largeur totale
# rangée haute (hauteur) et rangée basse (coûts/revenus)
TOP_H = Inches(4.7)
BOT_H = Inches(1.62)
GAP = Inches(0.07)
# 5 colonnes : KP | (KA/KR) | VP | (CR/CH) | CS
col_w = Inches(2.46)       # colonnes 1,3,5
mid_w = Inches(2.46)       # colonnes 2,4 (empilées)
# largeurs : KP, KA/KR, VP, CR/CH, CS — on fait 5 colonnes égales
cw = Emu(int((int(GW) - 4 * int(GAP)) / 5))
half_h = Emu(int((int(TOP_H) - int(GAP)) / 2))

x0 = int(GX)
xs = [Emu(x0 + i * (int(cw) + int(GAP))) for i in range(5)]
ty = int(GY)

# Bloc 1 — Partenaires clés (col 0, pleine hauteur)
bmc_block(s, xs[0], Emu(ty), cw, TOP_H, "07", "Partenaires clés", TEAL, [
    "Fournisseurs infra souveraine UE/Maroc (Scaleway, OVH)",
    "Hébergeurs self-host / on-prem",
    "Proxy LiteLLM (modèles interchangeables)",
    "Firecrawl (veille prix + catalogue)",
    "Open-source : embeddings Apache 2.0 (Qwen3-VL), LLM ouverts",
    "Intégrateurs / cabinets conseil (exécution Lot 2)",
    "Cabinets juridiques / DPO (couche conformité)",
    "Communauté d'Amine (P1, premiers clients)",
])
# Bloc 2 — Activités clés (col 1, haut)
bmc_block(s, xs[1], Emu(ty), cw, half_h, "06", "Activités clés", BLUE, [
    "Moteur de reco 7 couches (déterministe, pur)",
    "Price feed + veille catalogue live (sourcés)",
    "Génération Exit Escrow reproductible",
    "Déploiement assisté + supervision (Lot 2)",
], item_fs=8.6, sa=2.0)
# Bloc Ressources clés (col 1, bas)
bmc_block(s, xs[1], Emu(ty + int(half_h) + int(GAP)), cw, half_h, "08", "Ressources clés", BLUE, [
    "Codebase moteur pur + 3 moats (IP)",
    "Catalogue/baseline prix sourcés (DÉFCON 1)",
    "Dataset coûts réels (Network, à venir)",
    "Marque souveraine + charte fiduciaire",
], item_fs=8.6, sa=2.0)
# Bloc 3 — Proposition de valeur (col 2, pleine hauteur, accent)
rrect(s, xs[2], Emu(ty), cw, TOP_H, RGBColor(0xE6, 0xF5, 0xF2), border=TEAL, border_w=2.0, radius=0.03)
hd = rect(s, xs[2], Emu(ty), cw, Inches(0.34), TEAL)
htf = hd.text_frame; htf.word_wrap = True; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
htf.margin_left = Inches(0.1); htf.margin_top = 0; htf.margin_bottom = 0
hp = htf.paragraphs[0]
r1 = hp.add_run(); set_run(r1, "01  ", 11.5, WHITE, FONT_MONO, bold=True)
r2 = hp.add_run(); set_run(r2, "Proposition de valeur", 11, WHITE, FONT_TITLE, bold=True)
tb, tf = textbox(s, Emu(int(xs[2]) + int(Inches(0.1))), Emu(ty + int(Inches(0.42))),
                 Emu(int(cw) - int(Inches(0.2))), Emu(int(TOP_H) - int(Inches(0.5))))
para(tf, "« Base mémorielle souveraine qui grandit avec vous — sans migration, sans verrouillage, sans coûts cachés. »",
     10.5, TEAL, FONT_TITLE, bold=True, italic=True, first=True, sa=7, ls=1.0)
for it in [
    ("Souveraineté ", "(vos données, votre loi : UE/Maroc/on-prem)"),
    ("Anti-lock-in ", "(Exit Escrow : sortie en 1 clic, prouvée)"),
    ("Conseil honnête ", "(fiduciaire, zéro commission cachée)"),
    ("Coûts sourcés ±30 % ", "(URL+date+confiance, feed live)"),
    ("Décision sans expert ", "(verdict 90 s + intake libre)"),
]:
    p = tf.add_paragraph(); p.space_after = Pt(3.5); p.line_spacing = 1.0
    rm = p.add_run(); set_run(rm, "› ", 9.0, TEAL, FONT_BODY, bold=True)
    rb = p.add_run(); set_run(rb, it[0], 9.0, INK, FONT_BODY, bold=True)
    rt = p.add_run(); set_run(rt, it[1], 9.0, INK_VARIANT, FONT_BODY)
# Bloc 4 — Relations clients (col 3, haut)
bmc_block(s, xs[3], Emu(ty), cw, half_h, "04", "Relations clients", GOLD_TXT, [
    "Self-service (diagnostic 90 s gratuit)",
    "Assistant Q&A contextuel (sourcé)",
    "Mandat fiduciaire = confiance durable",
    "Accompagnement humain sur la cuisine",
], item_fs=8.6, sa=2.0)
# Bloc Canaux (col 3, bas)
bmc_block(s, xs[3], Emu(ty + int(half_h) + int(GAP)), cw, half_h, "03", "Canaux", GOLD_TXT, [
    "Site / configurateur web (PLG)",
    "Communauté d'Amine + bouche-à-oreille",
    "Contenu / démo live Exit Escrow",
    "Livrable exportable (MD/PDF) viral",
], item_fs=8.6, sa=2.0)
# Bloc 5 — Segments de clientèle (col 4, pleine hauteur)
bmc_block(s, xs[4], Emu(ty), cw, TOP_H, "02", "Segments de clientèle", GOLD_TXT, [
    ("P1 cœur ", "— bâtisseur souverain (coach/consultant FR Maroc-UE)"),
    ("P2 ", "— PME / startup tech (5-10 pers.)"),
    ("P3 à venir ", "— agence / programme (multi-tenant)"),
    ("P4 à venir ", "— cabinet régulé (secret pro, HARD)"),
    ("P5 hors périmètre ", "— ETI / grand groupe (appel d'offres)"),
    "Marché B2B, francophone d'abord",
])

# Rangée basse : Structure de coûts | Flux de revenus
by = ty + int(TOP_H) + int(GAP)
bot_w = Emu(int((int(GW) - int(GAP)) / 2))
# Structure de coûts
rrect(s, GX, Emu(by), bot_w, BOT_H, WHITE, border=CARD_BORDER, border_w=1.0, radius=0.03)
hd = rect(s, GX, Emu(by), bot_w, Inches(0.32), ERROR)
htf = hd.text_frame; htf.vertical_anchor = MSO_ANCHOR.MIDDLE; htf.margin_left = Inches(0.1)
hp = htf.paragraphs[0]
r1 = hp.add_run(); set_run(r1, "09  ", 11, WHITE, FONT_MONO, bold=True)
r2 = hp.add_run(); set_run(r2, "Structure de coûts", 11, WHITE, FONT_TITLE, bold=True)
tb, tf = textbox(s, Emu(int(GX) + int(Inches(0.12))), Emu(by + int(Inches(0.4))),
                 Emu(int(bot_w) - int(Inches(0.24))), Emu(int(BOT_H) - int(Inches(0.46))))
p = tf.paragraphs[0]
for seg in [("Coûts variables IA ", True), ("(LLM via LiteLLM + Firecrawl, même free tier) · ", False),
            ("Dev & maintenance ", True), ("(moteur, moats, veille) · ", False),
            ("Hébergement plateforme ", True), ("(self-host) · ", False),
            ("Support & CAC ", True), ("(≈0 sur P1) · ", False),
            ("infra client = REVERSÉE, hors marge Strate", False)]:
    r = p.add_run(); set_run(r, seg[0], 9.0, (INK if seg[1] else INK_VARIANT), FONT_BODY, bold=seg[1])
p.line_spacing = 1.05
para(tf, "Modèle orienté coûts variables IA à borner. Tout montant = hypothèse.", 8.5, ERROR, FONT_BODY, italic=True, sa=0, sb=3)
# Flux de revenus
rx = int(GX) + int(bot_w) + int(GAP)
rrect(s, Emu(rx), Emu(by), bot_w, BOT_H, WHITE, border=CARD_BORDER, border_w=1.0, radius=0.03)
hd = rect(s, Emu(rx), Emu(by), bot_w, Inches(0.32), TEAL)
htf = hd.text_frame; htf.vertical_anchor = MSO_ANCHOR.MIDDLE; htf.margin_left = Inches(0.1)
hp = htf.paragraphs[0]
r1 = hp.add_run(); set_run(r1, "05  ", 11, WHITE, FONT_MONO, bold=True)
r2 = hp.add_run(); set_run(r2, "Flux de revenus", 11, WHITE, FONT_TITLE, bold=True)
tb, tf = textbox(s, Emu(rx + int(Inches(0.12))), Emu(by + int(Inches(0.4))),
                 Emu(int(bot_w) - int(Inches(0.24))), Emu(int(BOT_H) - int(Inches(0.46))))
p = tf.paragraphs[0]
for seg in [("Diagnostic gratuit ", True), ("(lead magnet) · ", False),
            ("Pro ", True), ("= forfait service /mois [PLACEHOLDER] · ", False),
            ("Souverain+ ", True), ("sur devis · ", False),
            ("Négo vendor ", True), ("fiduciaire (rémunération affichée) · ", False),
            ("setup one-time + run récurrent", False)]:
    r = p.add_run(); set_run(r, seg[0], 9.0, (INK if seg[1] else INK_VARIANT), FONT_BODY, bold=seg[1])
p.line_spacing = 1.05
para(tf, "On tarife le SERVICE, jamais l'infra (zéro marge dessus). Prix à caler par le sondage.", 8.5, TEAL, FONT_BODY, italic=True, sa=0, sb=3)

# légende couleurs
tb, tf = textbox(s, Inches(0.35), Inches(7.16), Inches(12.6), Inches(0.3))
p = tf.paragraphs[0]
r = p.add_run(); set_run(r, "Lecture : la valeur (centre) sert les segments (droite) via relations & canaux ; elle est produite par activités/ressources/partenaires (gauche) ; coûts vs revenus en bas.", 9, INK_VARIANT, FONT_BODY, italic=True)

notes(s, """
LE CANEVAS EN PLÉNIÈRE. Présenter dans l'ordre de la logique Osterwalder : commencer par le CENTRE (proposition de valeur) et la DROITE (pour qui), puis remonter à GAUCHE (comment on le produit), enfin le BAS (combien ça coûte / combien ça rapporte).

PROPOSITION DE VALEUR (01) : « Base mémorielle souveraine qui grandit avec vous — sans migration, sans verrouillage, sans coûts cachés. » Cinq promesses : souveraineté (UE/Maroc/on-prem), anti-lock-in (Exit Escrow, seule preuve DURE), conseil honnête (fiduciaire), coûts sourcés ±30 % (feed live), décision sans expert (verdict 90 s + intake libre). On vend un RÉSULTAT souverain + une tranquillité, pas une stack.

SEGMENTS (02) : P1 (bâtisseur souverain FR Maroc-UE) = cœur/wedge (accès direct, CAC ≈0) ; P2 (PME tech) secondaire ; P3 (agences) et P4 (régulés) à venir ; P5 (ETI/grand groupe) hors périmètre self-service. B2B francophone d'abord.

CANAUX (03) : configurateur web (PLG, diagnostic 90 s gratuit), communauté + bouche-à-oreille, contenu/démo live de l'Exit Escrow (moment wow), livrable exportable MD/PDF (effet viral léger).

RELATIONS CLIENTS (04) : self-service en haut de funnel + assistant Q&A contextuel sourcé ; le mandat fiduciaire crée une confiance durable ; accompagnement humain réservé à la « cuisine » payante.

FLUX DE REVENUS (05) : Diagnostic gratuit (lead magnet) ; Pro = forfait service /mois = [PLACEHOLDER] (à caler par le sondage) ; Souverain+ sur devis ; négo vendor fiduciaire (rémunération affichée) ; setup one-time + run récurrent. On tarife le SERVICE, jamais l'infra.

ACTIVITÉS (06) : moteur de reco 7 couches déterministe, price feed + veille catalogue live sourcés, génération Exit Escrow, déploiement assisté + supervision (Lot 2).

PARTENAIRES (07) : fournisseurs infra souveraine (Scaleway/OVH), hébergeurs self-host, proxy LiteLLM, Firecrawl, open-source (embeddings Apache 2.0), intégrateurs/conseil, juridiques/DPO, communauté d'Amine.

RESSOURCES (08) : codebase moteur pur + 3 moats (IP), catalogue/baseline de prix sourcés, futur dataset de coûts réels (Network), marque souveraine + charte fiduciaire.

STRUCTURE DE COÛTS (09) : coûts VARIABLES IA (LLM + Firecrawl, même sur le free tier — à borner), dev/maintenance, hébergement plateforme, support, CAC (≈0 sur P1). L'infra du client est REVERSÉE sans marge → hors P&L de Strate. Modèle orienté coûts variables IA. Tout montant = hypothèse.

DÉFCON 1 : prix = [PLACEHOLDER], revenus/coûts = hypothèses, aucune stat inventée.
""")

# ===========================================================================
# SLIDE 3 — Détail blocs cœur (VP / Segments / Canaux / Relations)
# ===========================================================================
s = add_slide(SURFACE)
bar = rect(s, Inches(0), Inches(0), Inches(0.16), Inches(7.5), TEAL)
tb, tf = textbox(s, Inches(0.55), Inches(0.45), Inches(12.0), Inches(0.4))
p = para(tf, "DÉTAIL DES BLOCS · CÔTÉ MARCHÉ", 12.5, TEAL, FONT_BODY, bold=True, first=True)
for r in p.runs:
    r._r.get_or_add_rPr().set('spc', '140')
tb, tf = textbox(s, Inches(0.55), Inches(0.82), Inches(12.0), Inches(0.7))
para(tf, "Valeur · Segments · Canaux · Relations", 26, INK, FONT_TITLE, bold=True, first=True)
blocks = [
    ("01 · Proposition de valeur", TEAL, [
        "Job : décider, exécuter et maintenir une infra mémorielle engageante, sans expertise et sans capture.",
        "Preuve DURE = Exit Escrow + charte fiduciaire ; le reste « à mesurer chez vous » (zéro stat inventée).",
    ]),
    ("02 · Segments de clientèle", GOLD_TXT, [
        "P1 cœur (communauté d'Amine) → demande quasi acquise ; P2 extension technophile.",
        "P3/P4 à venir (multi-tenant déjà câblé / premium régulé) ; P5 explicitement hors périmètre.",
    ]),
    ("03 · Canaux", BLUE, [
        "PLG : diagnostic 90 s gratuit comme aimant (modèle Cloudflare free tier).",
        "Bouche-à-oreille communautaire + démonstration live de l'Exit Escrow (moment wow).",
    ]),
    ("04 · Relations clients", GOLD_TXT, [
        "Self-service + assistant Q&A sourcé en haut de funnel ; humain sur la cuisine payante.",
        "Le mandat fiduciaire (zéro commission) ancre une relation de confiance durable.",
    ]),
]
top = Inches(1.7); rh = Inches(1.22); gap = Inches(0.14); l = Inches(0.55); w = Inches(12.2)
for i, (h, col, items) in enumerate(blocks):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    rrect(s, l, t, w, rh, WHITE, border=CARD_BORDER, border_w=1.0, radius=0.04)
    sb = rect(s, l, t, Inches(0.08), rh, col)
    tb, tf = textbox(s, Emu(int(l) + int(Inches(0.28))), Emu(int(t) + int(Inches(0.13))),
                     Emu(int(w) - int(Inches(0.5))), Emu(int(rh) - int(Inches(0.24))))
    para(tf, h, 14, col, FONT_TITLE, bold=True, first=True, sa=5)
    for it in items:
        p = tf.add_paragraph(); p.space_after = Pt(2); p.line_spacing = 1.0
        rm = p.add_run(); set_run(rm, "— ", 11, col, FONT_BODY, bold=True)
        rt = p.add_run(); set_run(rt, it, 11, INK_VARIANT, FONT_BODY)
notes(s, """
Détail des 4 blocs « côté marché » (la moitié droite + centre du canevas).

01 PROPOSITION DE VALEUR — Le job (JTBD) : permettre à un non-expert de prendre, exécuter et maintenir dans le temps une décision d'infra engageante et difficilement réversible, sans se faire capturer. Distinction de régime de preuve essentielle (DÉFCON 1) : seules deux promesses sont démontrables aujourd'hui — l'Exit Escrow (le bundle redéploie une stack identique, c'est testé) et la charte fiduciaire (engagement contractuel). Le gain de temps / la qualité de réponse = « à mesurer chez vous », jamais chiffrés sans donnée.

02 SEGMENTS — La cible couvre 5 ordres de grandeur, donc la discipline ICP est cruciale. On commence par P1 (bâtisseur souverain de la communauté d'Amine) car la demande y est quasi acquise et le CAC ≈0. P2 (PME tech) suit. P3 (agences) et P4 (régulés) sont des lots ultérieurs : plus de valeur mais cycle plus lourd. P5 (ETI/grand groupe) est hors périmètre self-service (appel d'offres, DPA, SecNumCloud). Dire non à P5 est une décision, pas une faiblesse.

03 CANAUX — Product-Led Growth : le diagnostic gratuit (4 questions ou intake en langage libre) est le canal d'entrée. La communauté d'Amine + le bouche-à-oreille portent l'acquisition initiale. La démonstration live de l'Exit Escrow est le « moment wow » qui matérialise l'anti-lock-in. Le livrable exportable (MD/PDF) circule et fait connaître la marque.

04 RELATIONS CLIENTS — Self-service automatisé en haut de funnel (configurateur, assistant Q&A contextuel et sourcé), accompagnement humain réservé à la « cuisine » payante (déploiement, supervision). Le ressort de fidélité central est le MANDAT FIDUCIAIRE : être payé par le client et jamais commissionné par les vendors crée une relation de confiance qu'un acteur commissionné ne peut pas répliquer (leçon Flipper).
""")

# ===========================================================================
# SLIDE 4 — Détail blocs infra (Activités / Ressources / Partenaires)
# ===========================================================================
s = add_slide(SURFACE)
bar = rect(s, Inches(0), Inches(0), Inches(0.16), Inches(7.5), BLUE)
tb, tf = textbox(s, Inches(0.55), Inches(0.45), Inches(12.0), Inches(0.4))
p = para(tf, "DÉTAIL DES BLOCS · CÔTÉ PRODUCTION", 12.5, BLUE, FONT_BODY, bold=True, first=True)
for r in p.runs:
    r._r.get_or_add_rPr().set('spc', '140')
tb, tf = textbox(s, Inches(0.55), Inches(0.82), Inches(12.0), Inches(0.7))
para(tf, "Activités · Ressources · Partenaires", 26, INK, FONT_TITLE, bold=True, first=True)
blocks = [
    ("06 · Activités clés", BLUE, [
        "Moteur de reco 7 couches PUR et déterministe (le LLM propose, ne calcule jamais).",
        "Price feed + veille catalogue live (Firecrawl + LLM) réconciliés à une baseline sourcée.",
        "Génération de l'Exit Escrow reproductible ; déploiement assisté + supervision (Lot 2, à venir).",
    ]),
    ("08 · Ressources clés", BLUE, [
        "Codebase du moteur pur + les 3 moats (propriété intellectuelle).",
        "Catalogue + baseline de prix sourcés (DÉFCON 1) ; futur dataset de coûts réels (Network).",
        "Marque souveraine + charte fiduciaire (actif de confiance).",
    ]),
    ("07 · Partenaires clés", TEAL, [
        "Infra souveraine UE/Maroc (Scaleway, OVH), hébergeurs self-host/on-prem.",
        "Proxy LiteLLM (modèles interchangeables) + Firecrawl ; open-source (embeddings Apache 2.0).",
        "Intégrateurs/conseil (exécution Lot 2), juridiques/DPO (conformité), communauté d'Amine (P1).",
    ]),
]
top = Inches(1.75); rh = Inches(1.6); gap = Inches(0.16); l = Inches(0.55); w = Inches(12.2)
for i, (h, col, items) in enumerate(blocks):
    t = Emu(int(top) + i * (int(rh) + int(gap)))
    rrect(s, l, t, w, rh, WHITE, border=CARD_BORDER, border_w=1.0, radius=0.04)
    sb = rect(s, l, t, Inches(0.08), rh, col)
    tb, tf = textbox(s, Emu(int(l) + int(Inches(0.28))), Emu(int(t) + int(Inches(0.16))),
                     Emu(int(w) - int(Inches(0.5))), Emu(int(rh) - int(Inches(0.3))))
    para(tf, h, 14, col, FONT_TITLE, bold=True, first=True, sa=6)
    for it in items:
        p = tf.add_paragraph(); p.space_after = Pt(3); p.line_spacing = 1.0
        rm = p.add_run(); set_run(rm, "— ", 11, col, FONT_BODY, bold=True)
        rt = p.add_run(); set_run(rt, it, 11, INK_VARIANT, FONT_BODY)
notes(s, """
Détail des 3 blocs « côté production » (moitié gauche du canevas) : comment la valeur est créée.

06 ACTIVITÉS CLÉS — (a) Le moteur de recommandation 7 couches (C0→C6) est PUR et 100 % déterministe : le LLM propose des composants sourcés mais ne calcule JAMAIS un coût. C'est ce qui préserve la reproductibilité (Exit Escrow) et l'explicabilité fiduciaire. (b) Le price feed et la veille catalogue sont LIVE (Firecrawl + LLM via LiteLLM) mais chaque valeur passe un garde-fou : confrontée à une baseline sourcée, une valeur aberrante est rejetée au profit de la baseline (DÉFCON 1). (c) La génération de l'Exit Escrow (bundle reproductible). (d) Le déploiement assisté + la supervision = activité du Lot 2 (à venir), c'est là que naît le revenu récurrent.

08 RESSOURCES CLÉS — (a) La codebase du moteur pur et des 3 moats = la propriété intellectuelle défendable. (b) Le catalogue de composants + la baseline de prix sourcés (audit trail complet : figure + unité + URL + date + confiance). (c) Le futur dataset de coûts réels anonymisés (moat ③ Network) — ressource qui ne se constitue qu'avec des déploiements monitorés (Lot 2). (d) La marque souveraine + la charte fiduciaire, actifs immatériels de confiance.

07 PARTENAIRES CLÉS — (a) Fournisseurs d'infra souveraine (Scaleway, OVH) et hébergeurs self-host/on-prem : ce sont les briques que Strate recommande et, en Fiduciary Mode, négocie au nom du client (rémunération affichée, jamais de rétrocommission cachée). (b) Le proxy LiteLLM rend les modèles interchangeables (anti-dépendance, souveraineté) ; Firecrawl pour la veille ; open-source pour les embeddings (Apache 2.0, self-hostables). (c) Intégrateurs/cabinets pour exécuter le Lot 2 à l'échelle, juridiques/DPO pour la couche conformité (Strate reste « ingénierie, pas juridique »), et la communauté d'Amine comme premier vivier de clients.

Tension à rappeler : tout partenariat doit servir l'un des 3 moats, jamais les diluer — un partenariat commissionné qui orienterait la reco serait un anti-pattern à refuser, même lucratif.
""")

# ===========================================================================
# SLIDE 5 — Détail blocs financiers (Revenus / Coûts) + garde-fous
# ===========================================================================
s = add_slide(INVERSE)
bar = rect(s, Inches(0), Inches(0), Inches(0.16), Inches(7.5), GOLD)
tb, tf = textbox(s, Inches(0.55), Inches(0.45), Inches(12.0), Inches(0.4))
p = para(tf, "DÉTAIL DES BLOCS · ÉCONOMIE", 12.5, RGBColor(0xFF, 0xBA, 0x46), FONT_BODY, bold=True, first=True)
for r in p.runs:
    r._r.get_or_add_rPr().set('spc', '140')
tb, tf = textbox(s, Inches(0.55), Inches(0.82), Inches(12.0), Inches(0.7))
para(tf, "Flux de revenus · Structure de coûts", 26, WHITE, FONT_TITLE, bold=True, first=True)
# deux colonnes
c1 = rrect(s, Inches(0.55), Inches(1.75), Inches(6.0), Inches(4.5), RGBColor(0x32, 0x3b, 0x52), border=None, radius=0.04)
tb, tf = textbox(s, Inches(0.78), Inches(1.95), Inches(5.55), Inches(4.2))
para(tf, "05 · Flux de revenus", 15, MINT, FONT_TITLE, bold=True, first=True, sa=9)
for it in [
    ("Diagnostic ", "= gratuit (lead magnet, crée l'habitude)."),
    ("Pro ", "= forfait service /mois = [PLACEHOLDER]."),
    ("Souverain+ ", "= sur devis (on-prem, SLA, multi-région à venir)."),
    ("Négo vendor ", "= rémunération fiduciaire AFFICHÉE."),
    ("Mise en route ", "= setup one-time + run récurrent."),
]:
    p = tf.add_paragraph(); p.space_after = Pt(7); p.line_spacing = 1.05
    rm = p.add_run(); set_run(rm, "› ", 12, MINT, FONT_BODY, bold=True)
    rb = p.add_run(); set_run(rb, it[0], 12, WHITE, FONT_BODY, bold=True)
    rt = p.add_run(); set_run(rt, it[1], 12, INVERSE_TXT, FONT_BODY)
para(tf, "On tarife le SERVICE ; l'infra reste au coût fournisseur (zéro marge cachée).", 11, MINT, FONT_BODY, italic=True, sb=4)
c2 = rrect(s, Inches(6.75), Inches(1.75), Inches(6.0), Inches(4.5), RGBColor(0x32, 0x3b, 0x52), border=None, radius=0.04)
tb, tf = textbox(s, Inches(6.98), Inches(1.95), Inches(5.55), Inches(4.2))
para(tf, "09 · Structure de coûts", 15, RGBColor(0xFF, 0xBA, 0x46), FONT_TITLE, bold=True, first=True, sa=9)
for it in [
    ("Coûts variables IA ", "(LLM via LiteLLM + Firecrawl), même au free tier."),
    ("Dev & maintenance ", "(moteur, moats, veille sourcée)."),
    ("Hébergement plateforme ", "(self-host Coolify/serveur)."),
    ("Support & CAC ", "(≈0 sur P1, croissant ensuite)."),
    ("Infra client ", "= REVERSÉE sans marge → hors P&L Strate."),
]:
    p = tf.add_paragraph(); p.space_after = Pt(7); p.line_spacing = 1.05
    rm = p.add_run(); set_run(rm, "› ", 12, RGBColor(0xFF, 0xBA, 0x46), FONT_BODY, bold=True)
    rb = p.add_run(); set_run(rb, it[0], 12, WHITE, FONT_BODY, bold=True)
    rt = p.add_run(); set_run(rt, it[1], 12, INVERSE_TXT, FONT_BODY)
para(tf, "Modèle orienté coûts variables IA → à borner (quotas/cache).", 11, RGBColor(0xFF, 0xBA, 0x46), FONT_BODY, italic=True, sb=4)
# bandeau garde-fou
tb, tf = textbox(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.7))
para(tf, "DÉFCON 1 : prix de vente = [PLACEHOLDER] (sondage Van Westendorp + conjoint en cours) · tout montant de revenus/coûts = hypothèse à valider · aucune stat de marché inventée.",
     11.5, RGBColor(0xFF, 0xB4, 0xAB), FONT_MONO, italic=False, first=True, ls=1.05)
notes(s, """
Détail des 2 blocs financiers + rappel des garde-fous. Fond sombre pour marquer qu'on entre dans l'économie HYPOTHÉTIQUE.

05 FLUX DE REVENUS — Architecture « recette ouverte, cuisine payante » :
- Diagnostic = GRATUIT (lead magnet, modèle Cloudflare free tier). Le QUOI est offert (anti-lock-in + habitude).
- Pro = forfait service mensuel = [PLACEHOLDER]. C'est LE prix à caler par le sondage. Inclut déploiement assisté + supervision + recalibration + négo vendor + Exit bundle premium.
- Souverain+ = sur devis (on-prem/air-gapped, gouvernance, SLA, multi-région à venir). ARPA supérieur, cycle plus long.
- Négo vendor : rémunération fiduciaire AFFICHÉE (jamais une rétrocommission cachée) — compatible avec le moat ② à condition de divulgation.
- Mise en route : setup one-time (ingestion du backlog) + run récurrent.
Principe : on tarife le SERVICE, jamais l'infra. Margerait-on sur l'infra qu'on aurait intérêt à pousser des composants chers — exactement le conflit qu'on dénonce.

09 STRUCTURE DE COÛTS — Modèle plutôt « orienté coûts variables IA » :
- Coûts variables IA (LLM via LiteLLM + Firecrawl) : déclenchés même par un utilisateur GRATUIT (intake, narration, assistant, veille). C'est l'angle aveugle économique principal → à borner par quotas/cache/modèle économe.
- Dev & maintenance (moteur, moats, veille sourcée).
- Hébergement plateforme (self-host Coolify/serveur, léger au démarrage).
- Support & CAC (≈0 sur P1, croissant hors réseau).
- L'infra du CLIENT est REVERSÉE sans marge → elle n'entre PAS dans le P&L de Strate. Donc le revenu Strate est du service à forte marge brute logicielle, MAIS nette des coûts variables IA.

GARDE-FOUS DÉFCON 1 (à dire à voix haute) : le prix est un [PLACEHOLDER] tant que le sondage (≥15 Van Westendorp + ≥30 conjoint) n'a pas ses seuils ; tout chiffre de revenus/coûts est une hypothèse de modélisation ; aucune taille de marché n'est inventée.

Le vrai enjeu de viabilité n'est pas « marge par client » (le service marge bien) mais « le coût IA du free tier est-il amorti par la conversion ? » et « le CAC hors-communauté est-il soutenable ? ». D'où la priorité : clore le sondage + réaliser 5 ventes P1 pour observer conversion, ARPA effectif, churn et coût IA/lead.
""")

# ---------------------------------------------------------------------------
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "strate-business-model-canvas.pptx")
prs.save(OUT)
print("OK ->", OUT)
print("slides:", len(prs.slides._sldIdLst))
